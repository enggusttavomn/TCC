"""Gera uma apresentação de 10 slides sobre os modelos e resultados do TCC."""

from __future__ import annotations

import csv
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


RAIZ = Path(__file__).resolve().parents[1]
SAIDA = RAIZ / "00_APRESENTACAO" / "modelos_baselines_resultados_10_slides.pptx"
RESULTADOS = RAIZ / "resultados" / "avaliacao_mensal_canonica"

BG = "071724"
CARD = "102B3D"
CARD_ALT = "0D2232"
GRID = "25465A"
WHITE = "F5F8FA"
MUTED = "A9BAC5"
TEAL = "2DD4BF"
ORANGE = "FF825C"
BLUE = "55A7FF"
PURPLE = "8B7CF6"
LIME = "C5F55A"
PINK = "F46AA5"
YELLOW = "F7C85B"
GREEN = "39B87F"
RED = "FF6B6B"

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)


def rgb(hexadecimal: str) -> RGBColor:
    return RGBColor.from_string(hexadecimal)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color: str = WHITE,
    bold: bool = False,
    font: str = "Aptos",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.04,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_rich_text(slide, runs, x, y, w, h, *, size=16, color=WHITE):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.06)
    frame.margin_bottom = Inches(0.06)
    p = frame.paragraphs[0]
    for item in runs:
        run = p.add_run()
        run.text = item["text"]
        run.font.name = "Aptos"
        run.font.size = Pt(item.get("size", size))
        run.font.bold = item.get("bold", False)
        run.font.color.rgb = rgb(item.get("color", color))
    return box


def add_shape(
    slide,
    shape_type,
    x,
    y,
    w,
    h,
    *,
    fill=CARD,
    line=None,
    radius=True,
    transparency=0,
):
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = transparency
    shape.line.color.rgb = rgb(line or fill)
    shape.line.width = Pt(0.8)
    return shape


def add_line(slide, x1, y1, x2, y2, *, color=GRID, width=1.3, dash=None):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    if dash is not None:
        line.line.dash_style = dash
    return line


def add_circle(slide, x, y, d, *, fill, line=None):
    return add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.OVAL,
        x,
        y,
        d,
        d,
        fill=fill,
        line=line or fill,
    )


def add_tag(slide, text, x, y, *, color=TEAL, width=1.65):
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        x,
        y,
        width,
        0.32,
        fill=CARD_ALT,
        line=color,
    )
    add_text(
        slide,
        text.upper(),
        x + 0.06,
        y + 0.02,
        width - 0.12,
        0.23,
        size=8.5,
        color=color,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0,
    )


def add_header(slide, category, title, subtitle, number, accent):
    add_text(
        slide,
        category.upper(),
        0.7,
        0.28,
        4.0,
        0.24,
        size=9,
        color=accent,
        bold=True,
        margin=0,
    )
    add_text(
        slide,
        title,
        0.7,
        0.63,
        10.9,
        0.52,
        size=27,
        bold=True,
        margin=0,
    )
    add_text(
        slide,
        subtitle,
        0.7,
        1.18,
        11.4,
        0.36,
        size=12.5,
        color=MUTED,
        margin=0,
    )
    add_text(
        slide,
        f"{number:02d}",
        12.08,
        0.29,
        0.5,
        0.24,
        size=9,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
        margin=0,
    )
    add_line(slide, 0.7, 1.57, 12.63, 1.57, color=GRID, width=0.8)


def add_footer(slide, label):
    add_text(
        slide,
        "TCC · PREVISÃO MENSAL DE GHI · 10 LOCALIDADES · TESTE EM 2024",
        0.7,
        7.16,
        6.8,
        0.16,
        size=6.5,
        color="658092",
        margin=0,
    )
    add_text(
        slide,
        label.upper(),
        10.5,
        7.16,
        2.13,
        0.16,
        size=6.5,
        color="658092",
        align=PP_ALIGN.RIGHT,
        margin=0,
    )


def add_metric_bar(slide, rank, mae, rmse, accent, note):
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        0.7,
        6.37,
        11.93,
        0.58,
        fill=CARD_ALT,
        line=GRID,
    )
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0.7,
        6.37,
        0.08,
        0.58,
        fill=accent,
        line=accent,
    )
    add_rich_text(
        slide,
        [
            {"text": f"MAE {mae:.2f} W/m²", "bold": True, "color": WHITE},
            {"text": f"   ·   RMSE {rmse:.2f} W/m²", "color": MUTED},
            {"text": f"   ·   {rank}º no ranking geral", "color": accent, "bold": True},
        ],
        0.93,
        6.49,
        5.7,
        0.28,
        size=12,
    )
    add_text(
        slide,
        note,
        6.6,
        6.47,
        5.75,
        0.28,
        size=10.5,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0,
    )


def add_info_card(slide, x, y, w, h, label, body, accent, *, body_size=13.2):
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        x,
        y,
        w,
        h,
        fill=CARD,
        line=GRID,
    )
    add_text(
        slide,
        label.upper(),
        x + 0.22,
        y + 0.18,
        w - 0.44,
        0.22,
        size=8.5,
        color=accent,
        bold=True,
        margin=0,
    )
    add_text(
        slide,
        body,
        x + 0.22,
        y + 0.52,
        w - 0.44,
        h - 0.68,
        size=body_size,
        color=WHITE,
        margin=0,
    )


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = rgb(BG)
    return slide


def draw_tree(slide, x, y, scale, accent):
    nodes = [
        (0.0, 0.8),
        (0.65, 0.32),
        (0.65, 1.28),
        (1.3, 0.05),
        (1.3, 0.6),
        (1.3, 1.0),
        (1.3, 1.55),
    ]
    edges = [(0, 1), (0, 2), (1, 3), (1, 4), (2, 5), (2, 6)]
    for a, b in edges:
        add_line(
            slide,
            x + nodes[a][0] * scale,
            y + nodes[a][1] * scale,
            x + nodes[b][0] * scale,
            y + nodes[b][1] * scale,
            color=GRID,
            width=1.1,
        )
    for idx, (nx, ny) in enumerate(nodes):
        add_circle(
            slide,
            x + nx * scale - 0.08,
            y + ny * scale - 0.08,
            0.16,
            fill=accent if idx in (0, 3, 6) else CARD_ALT,
            line=accent,
        )


def slide_xgboost(prs, metrics):
    slide = blank_slide(prs)
    accent = ORANGE
    add_header(
        slide,
        "Modelo tabular",
        "XGBoost: árvores que corrigem os erros anteriores",
        "Boosting combina várias árvores pequenas em uma previsão pontual.",
        1,
        accent,
    )
    add_info_card(
        slide,
        0.7,
        1.86,
        3.15,
        1.63,
        "Entrada",
        "12 lags mensais + médias de 3, 6 e 12 meses + calendário + localidade.",
        accent,
        body_size=13,
    )
    add_info_card(
        slide,
        0.7,
        3.69,
        3.15,
        1.83,
        "Treinamento",
        "Modelo global; validação temporal escolheu entre 207 e 392 árvores nas cinco sementes.",
        accent,
        body_size=12.6,
    )
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        4.12,
        1.86,
        5.35,
        3.66,
        fill=CARD,
        line=GRID,
    )
    add_tag(slide, "boosting", 4.38, 2.12, color=accent, width=1.15)
    for i, dx in enumerate((0.0, 1.25, 2.5)):
        draw_tree(slide, 4.68 + dx, 2.82, 0.72, accent if i == 2 else BLUE)
        if i < 2:
            add_text(
                slide,
                "+",
                5.83 + dx,
                3.27,
                0.3,
                0.3,
                size=20,
                color=MUTED,
                bold=True,
                align=PP_ALIGN.CENTER,
                margin=0,
            )
    add_line(slide, 8.45, 3.55, 8.95, 3.55, color=accent, width=2.2)
    add_circle(slide, 8.93, 3.37, 0.36, fill=accent)
    add_text(
        slide,
        "ŷ",
        8.94,
        3.39,
        0.34,
        0.3,
        size=14,
        color=BG,
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    add_text(
        slide,
        "árvores sucessivas  →  soma das correções",
        5.0,
        4.84,
        3.7,
        0.28,
        size=10,
        color=MUTED,
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    add_info_card(
        slide,
        9.72,
        1.86,
        2.91,
        1.63,
        "Ponto forte",
        "Excelente em atributos tabulares e interações não lineares.",
        accent,
        body_size=12.6,
    )
    add_info_card(
        slide,
        9.72,
        3.69,
        2.91,
        1.83,
        "Limite",
        "O tempo precisa ser explicitado por lags e variáveis derivadas.",
        accent,
        body_size=12.6,
    )
    m = metrics["XGBoost"]
    add_metric_bar(slide, 6, m["mae"], m["rmse"], accent, "Desempenho intermediário no teste")
    add_footer(slide, "XGBoost")


def slide_mlp(prs, metrics):
    slide = blank_slide(prs)
    accent = BLUE
    add_header(
        slide,
        "Modelo tabular",
        "MLP: camadas densas para relações não lineares",
        "Uma rede feedforward transforma os mesmos atributos tabulares em GHI mensal.",
        2,
        accent,
    )
    add_info_card(
        slide,
        0.7,
        1.86,
        3.15,
        1.63,
        "Arquitetura",
        "Duas camadas ocultas: 64 e 32 neurônios, ativação ReLU.",
        accent,
    )
    add_info_card(
        slide,
        0.7,
        3.69,
        3.15,
        1.83,
        "Regularização",
        "StandardScaler + penalização L2 (α = 0,05) e otimizador L-BFGS.",
        accent,
        body_size=12.8,
    )
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        4.12,
        1.86,
        5.35,
        3.66,
        fill=CARD,
        line=GRID,
    )
    columns = [
        (4.65, 5, MUTED),
        (5.88, 7, BLUE),
        (7.2, 5, TEAL),
        (8.55, 1, LIME),
    ]
    positions = {}
    for ci, (cx, n, color) in enumerate(columns):
        ys = [2.42 + j * (2.38 / max(1, n - 1)) for j in range(n)]
        positions[ci] = [(cx, yy) for yy in ys]
        for cxn, yy in positions[ci]:
            add_circle(slide, cxn, yy, 0.18, fill=color)
    for ci in range(len(columns) - 1):
        for x1, y1 in positions[ci]:
            for x2, y2 in positions[ci + 1]:
                add_line(
                    slide,
                    x1 + 0.18,
                    y1 + 0.09,
                    x2,
                    y2 + 0.09,
                    color="315166",
                    width=0.45,
                )
    for x, label in ((4.38, "atributos"), (5.67, "64"), (7.04, "32"), (8.39, "saída")):
        add_text(
            slide,
            label,
            x,
            5.03,
            0.65,
            0.22,
            size=8,
            color=MUTED,
            align=PP_ALIGN.CENTER,
            margin=0,
        )
    add_info_card(
        slide,
        9.72,
        1.86,
        2.91,
        1.63,
        "Ponto forte",
        "Flexível para interações suaves e complexas.",
        accent,
        body_size=12.8,
    )
    add_info_card(
        slide,
        9.72,
        3.69,
        2.91,
        1.83,
        "Limite",
        "Base pequena aumenta a sensibilidade à escala e à regularização.",
        accent,
        body_size=12.4,
    )
    m = metrics["MLP"]
    add_metric_bar(slide, 7, m["mae"], m["rmse"], accent, "Ficou abaixo dos modelos recorrentes")
    add_footer(slide, "MLP")


def draw_recurrent_chain(slide, accent, *, gates=False):
    xs = [4.55, 5.55, 6.55, 7.55, 8.55]
    for idx, x in enumerate(xs):
        add_shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            x,
            3.15,
            0.68,
            0.62,
            fill=CARD_ALT,
            line=accent,
        )
        add_text(
            slide,
            "LSTM" if gates else "RNN",
            x,
            3.34,
            0.68,
            0.2,
            size=8.5,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
            margin=0,
        )
        add_text(
            slide,
            f"xₜ₋{4-idx}" if idx < 4 else "xₜ",
            x + 0.11,
            2.63,
            0.46,
            0.2,
            size=10,
            color=accent,
            bold=True,
            align=PP_ALIGN.CENTER,
            margin=0,
        )
        add_line(slide, x + 0.34, 2.88, x + 0.34, 3.15, color=accent, width=1.1)
        if idx < len(xs) - 1:
            add_line(slide, x + 0.68, 3.46, xs[idx + 1], 3.46, color=accent, width=1.7)
    add_line(slide, 8.89, 3.77, 8.89, 4.2, color=LIME, width=1.7)
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        8.5,
        4.2,
        0.78,
        0.42,
        fill=LIME,
        line=LIME,
    )
    add_text(
        slide,
        "ŷₜ₊₁",
        8.5,
        4.31,
        0.78,
        0.18,
        size=10,
        color=BG,
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0,
    )


def slide_rnn(prs, metrics):
    slide = blank_slide(prs)
    accent = TEAL
    add_header(
        slide,
        "Modelo recorrente",
        "RNN: o estado oculto carrega o contexto",
        "A ordem dos 12 meses entra diretamente na arquitetura, passo a passo.",
        3,
        accent,
    )
    add_info_card(
        slide,
        0.7,
        1.86,
        3.15,
        1.63,
        "Memória",
        "Estado oculto de 16 unidades resume o histórico mais recente.",
        accent,
    )
    add_info_card(
        slide,
        0.7,
        3.69,
        3.15,
        1.83,
        "Fusão",
        "Calendário e localidade entram em um ramo auxiliar, após a recorrência.",
        accent,
        body_size=12.8,
    )
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        4.12,
        1.86,
        5.35,
        3.66,
        fill=CARD,
        line=GRID,
    )
    draw_recurrent_chain(slide, accent)
    add_text(
        slide,
        "hₜ = tanh(Wₓxₜ + Wₕhₜ₋₁ + b)",
        5.18,
        4.92,
        3.2,
        0.25,
        size=11,
        color=MUTED,
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    add_info_card(
        slide,
        9.72,
        1.86,
        2.91,
        1.63,
        "Ponto forte",
        "Representa naturalmente ordem e dependências locais.",
        accent,
        body_size=12.6,
    )
    add_info_card(
        slide,
        9.72,
        3.69,
        2.91,
        1.83,
        "Limite",
        "Sinais distantes podem enfraquecer durante o treinamento.",
        accent,
        body_size=12.6,
    )
    m = metrics["RNN"]
    add_metric_bar(slide, 3, m["mae"], m["rmse"], accent, "2º melhor modelo aprendido")
    add_footer(slide, "RNN")


def slide_lstm(prs, metrics):
    slide = blank_slide(prs)
    accent = LIME
    add_header(
        slide,
        "Modelo recorrente",
        "LSTM: memória seletiva para dependências longas",
        "Portas aprendem o que esquecer, adicionar e expor em cada mês.",
        4,
        accent,
    )
    add_info_card(
        slide,
        0.7,
        1.86,
        3.15,
        1.63,
        "Arquitetura",
        "16 unidades LSTM + fusão das covariáveis + camada densa.",
        accent,
    )
    add_info_card(
        slide,
        0.7,
        3.69,
        3.15,
        1.83,
        "Validação",
        "Épocas selecionadas temporalmente; depois, reajuste em todo o pré-teste.",
        accent,
        body_size=12.8,
    )
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        4.12,
        1.86,
        5.35,
        3.66,
        fill=CARD,
        line=GRID,
    )
    add_line(slide, 4.62, 2.62, 8.98, 2.62, color=accent, width=2.2)
    add_text(
        slide,
        "estado da célula cₜ",
        4.66,
        2.25,
        1.55,
        0.22,
        size=9,
        color=accent,
        bold=True,
        margin=0,
    )
    gate_data = [
        (5.05, "×", "ESQUECER", ORANGE),
        (6.45, "+", "ADICIONAR", TEAL),
        (7.85, "σ", "EXPOR", PURPLE),
    ]
    for gx, symbol, label, color in gate_data:
        add_circle(slide, gx, 3.1, 0.54, fill=CARD_ALT, line=color)
        add_text(
            slide,
            symbol,
            gx,
            3.24,
            0.54,
            0.2,
            size=15,
            color=color,
            bold=True,
            align=PP_ALIGN.CENTER,
            margin=0,
        )
        add_line(slide, gx + 0.27, 2.62, gx + 0.27, 3.1, color=color, width=1.25)
        add_text(
            slide,
            label,
            gx - 0.16,
            3.84,
            0.86,
            0.2,
            size=8,
            color=color,
            bold=True,
            align=PP_ALIGN.CENTER,
            margin=0,
        )
    add_text(
        slide,
        "xₜ + hₜ₋₁",
        6.22,
        4.48,
        1.2,
        0.26,
        size=12,
        color=MUTED,
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    add_line(slide, 6.82, 4.3, 6.82, 3.64, color=MUTED, width=1)
    add_info_card(
        slide,
        9.72,
        1.86,
        2.91,
        1.63,
        "Ponto forte",
        "Preserva sinais relevantes por horizontes maiores.",
        accent,
        body_size=12.6,
    )
    add_info_card(
        slide,
        9.72,
        3.69,
        2.91,
        1.83,
        "Limite",
        "Mais parâmetros e treinamento mais cuidadoso.",
        accent,
        body_size=12.6,
    )
    m = metrics["LSTM"]
    add_metric_bar(slide, 2, m["mae"], m["rmse"], accent, "Melhor modelo aprendido; quase empatou com a climatologia")
    add_footer(slide, "LSTM")


def slide_dilated(prs, metrics):
    slide = blank_slide(prs)
    accent = PURPLE
    add_header(
        slide,
        "Modelo recorrente",
        "DilatedRNN: conexões que saltam no tempo",
        "Dilatações 1, 2 e 4 constroem uma hierarquia temporal multiescala.",
        5,
        accent,
    )
    add_info_card(
        slide,
        0.7,
        1.86,
        3.15,
        1.63,
        "Estrutura",
        "Três camadas recorrentes com saltos s = 1, 2 e 4.",
        accent,
    )
    add_info_card(
        slide,
        0.7,
        3.69,
        3.15,
        1.83,
        "Intuição",
        "Capturar variações adjacentes e ritmos mais amplos sem aumentar a janela.",
        accent,
        body_size=12.8,
    )
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        4.12,
        1.86,
        5.35,
        3.66,
        fill=CARD,
        line=GRID,
    )
    xs = [4.55 + i * 0.38 for i in range(12)]
    for i, x in enumerate(xs):
        add_circle(slide, x, 4.62, 0.13, fill=CARD_ALT, line=accent if i in (0, 4, 8, 11) else GRID)
        add_text(
            slide,
            str(i),
            x - 0.03,
            4.91,
            0.2,
            0.15,
            size=6.5,
            color=MUTED,
            align=PP_ALIGN.CENTER,
            margin=0,
        )
    levels = [(1, 4.08, TEAL, "s = 1"), (2, 3.42, ORANGE, "s = 2"), (4, 2.62, PURPLE, "s = 4")]
    for step, yy, color, label in levels:
        add_tag(slide, label, 4.46, yy - 0.28, color=color, width=0.82)
        for start in range(0, 12 - step, step):
            x1 = xs[start] + 0.065
            x2 = xs[start + step] + 0.065
            add_line(slide, x1, yy, x2, yy, color=color, width=1.4)
            add_line(slide, x1, yy, x1, 4.62, color=color, width=0.45)
            add_line(slide, x2, yy, x2, 4.62, color=color, width=0.45)
    add_info_card(
        slide,
        9.72,
        1.86,
        2.91,
        1.63,
        "Ponto forte",
        "Acessa padrões em escalas temporais diferentes.",
        accent,
        body_size=12.6,
    )
    add_info_card(
        slide,
        9.72,
        3.69,
        2.91,
        1.83,
        "Limite",
        "O desenho das dilatações precisa combinar com o fenômeno.",
        accent,
        body_size=12.4,
    )
    m = metrics["DilatedRNN"]
    add_metric_bar(slide, 4, m["mae"], m["rmse"], accent, "Bom desempenho, mas abaixo de RNN e LSTM")
    add_footer(slide, "DilatedRNN")


def slide_deepar(prs, metrics, prob_metrics):
    slide = blank_slide(prs)
    accent = PINK
    add_header(
        slide,
        "Modelo probabilístico",
        "DeepAR: uma distribuição de futuros possíveis",
        "Uma rede autoregressiva global aprende padrões compartilhados entre as dez séries.",
        6,
        accent,
    )
    add_info_card(
        slide,
        0.7,
        1.86,
        3.15,
        1.63,
        "Configuração",
        "Contexto de 12 meses, 2 camadas, 40 unidades e 100 épocas.",
        accent,
    )
    add_info_card(
        slide,
        0.7,
        3.69,
        3.15,
        1.83,
        "Saída",
        "500 amostras por semente; 2.500 na mistura final. A mediana gera a previsão pontual.",
        accent,
        body_size=12.4,
    )
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        4.12,
        1.86,
        5.35,
        3.66,
        fill=CARD,
        line=GRID,
    )
    for row, color in enumerate((BLUE, TEAL, PURPLE, ORANGE)):
        points = []
        for i in range(7):
            px = 4.5 + i * 0.34
            py = 2.45 + row * 0.52 + ((i * (row + 2)) % 3) * 0.07
            points.append((px, py))
        for i in range(len(points) - 1):
            add_line(slide, *points[i], *points[i + 1], color=color, width=1.1)
    add_circle(slide, 7.15, 3.05, 0.7, fill=CARD_ALT, line=accent)
    add_text(
        slide,
        "RNN\nglobal",
        7.15,
        3.19,
        0.7,
        0.34,
        size=9,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    add_line(slide, 6.75, 3.4, 7.15, 3.4, color=accent, width=1.8)
    add_line(slide, 7.85, 3.4, 8.25, 3.4, color=accent, width=1.8)
    for idx, (label, yy, color) in enumerate((("P10", 2.62, PURPLE), ("P50", 3.28, PINK), ("P90", 3.94, YELLOW))):
        add_line(slide, 8.25, 3.4, 8.62, yy + 0.13, color=color, width=1.2)
        add_shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            8.62,
            yy,
            0.55,
            0.28,
            fill=color,
            line=color,
        )
        add_text(
            slide,
            label,
            8.62,
            yy + 0.065,
            0.55,
            0.14,
            size=7.5,
            color=BG,
            bold=True,
            align=PP_ALIGN.CENTER,
            margin=0,
        )
    p = prob_metrics["DeepAR"]
    add_info_card(
        slide,
        9.72,
        1.86,
        2.91,
        1.63,
        "Qualidade probabilística",
        f"CRPS {p['crps']:.2f} W/m²\nCobertura 90%: {p['picp']:.2f}%",
        accent,
        body_size=12.5,
    )
    add_info_card(
        slide,
        9.72,
        3.69,
        2.91,
        1.83,
        "Leitura",
        f"Intervalos relativamente estreitos ({p['mpiw']:.2f} W/m²), porém subcobertos.",
        accent,
        body_size=12.4,
    )
    m = metrics["DeepAR"]
    add_metric_bar(slide, 5, m["mae"], m["rmse"], accent, "Melhor qualidade probabilística entre os dois")
    add_footer(slide, "DeepAR")


def slide_deepnpts(prs, metrics, prob_metrics):
    slide = blank_slide(prs)
    accent = YELLOW
    add_header(
        slide,
        "Modelo probabilístico",
        "DeepNPTS: o histórico vira um conjunto de candidatos",
        "O estimador discreto oficial atribui massa de probabilidade aos valores observados no contexto.",
        7,
        accent,
    )
    add_info_card(
        slide,
        0.7,
        1.86,
        3.15,
        1.63,
        "Configuração",
        "DeepNPTSEstimator do GluonTS 0.16.2, versão discreta e treinamento global.",
        accent,
        body_size=12.5,
    )
    add_info_card(
        slide,
        0.7,
        3.69,
        3.15,
        1.83,
        "Correção técnica",
        "Embeddings categóricos registrados em ModuleList; arquitetura e perda RPS preservadas.",
        accent,
        body_size=12.2,
    )
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        4.12,
        1.86,
        5.35,
        3.66,
        fill=CARD,
        line=GRID,
    )
    stages = [
        (4.39, "1", "CONTEXTO", TEAL),
        (5.62, "2", "SCORE", ORANGE),
        (6.85, "3", "PESOS", PURPLE),
        (8.08, "4", "AMOSTRAS", YELLOW),
    ]
    for idx, (x, number, label, color) in enumerate(stages):
        add_shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            x,
            2.22,
            1.05,
            2.72,
            fill=CARD_ALT,
            line=GRID,
        )
        add_circle(slide, x + 0.1, 2.36, 0.24, fill=color)
        add_text(
            slide,
            number,
            x + 0.1,
            2.415,
            0.24,
            0.12,
            size=7,
            color=BG,
            bold=True,
            align=PP_ALIGN.CENTER,
            margin=0,
        )
        add_text(
            slide,
            label,
            x + 0.39,
            2.4,
            0.55,
            0.14,
            size=7.4,
            color=color,
            bold=True,
            margin=0,
        )
        if idx < len(stages) - 1:
            add_line(slide, x + 1.05, 3.58, stages[idx + 1][0], 3.58, color=color, width=1.6)
    for i in range(6):
        add_circle(slide, 4.54 + i * 0.13, 3.45 + (i % 2) * 0.15, 0.07, fill=TEAL)
        add_shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            5.82,
            3.1 + i * 0.25,
            0.2 + (i % 3) * 0.12,
            0.09,
            fill=ORANGE,
            line=ORANGE,
        )
    for i, d in enumerate((0.08, 0.12, 0.16, 0.22, 0.28)):
        add_circle(slide, 7.0 + i * 0.17, 3.46 - d / 2, d, fill=PURPLE)
    for i, (px, py, d) in enumerate(((8.3, 4.25, 0.12), (8.58, 3.63, 0.18), (8.86, 3.08, 0.25), (8.41, 2.95, 0.1))):
        add_circle(slide, px, py, d, fill=YELLOW)
    p = prob_metrics["DeepNPTS"]
    add_info_card(
        slide,
        9.72,
        1.86,
        2.91,
        1.63,
        "Qualidade probabilística",
        f"CRPS {p['crps']:.2f} W/m²\nCobertura 90%: {p['picp']:.2f}%",
        accent,
        body_size=12.5,
    )
    add_info_card(
        slide,
        9.72,
        3.69,
        2.91,
        1.83,
        "Leitura",
        f"Cobertura próxima do nominal, mas intervalos muito largos ({p['mpiw']:.2f} W/m²).",
        accent,
        body_size=12.1,
    )
    m = metrics["DeepNPTS"]
    add_metric_bar(slide, 9, m["mae"], m["rmse"], accent, "Alta variabilidade entre sementes (DP = 9,22 W/m²)")
    add_footer(slide, "DeepNPTS")


def baseline_card(slide, x, title, definition, result, rank, accent, symbol):
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        x,
        2.0,
        3.74,
        3.86,
        fill=CARD,
        line=GRID,
    )
    add_circle(slide, x + 0.25, 2.26, 0.45, fill=accent)
    add_text(
        slide,
        symbol,
        x + 0.25,
        2.36,
        0.45,
        0.2,
        size=13,
        color=BG,
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    add_text(
        slide,
        title,
        x + 0.25,
        2.88,
        3.24,
        0.38,
        size=18,
        color=WHITE,
        bold=True,
        margin=0,
    )
    add_text(
        slide,
        definition,
        x + 0.25,
        3.43,
        3.24,
        1.08,
        size=13.2,
        color=MUTED,
        margin=0,
    )
    add_line(slide, x + 0.25, 4.72, x + 3.49, 4.72, color=GRID, width=0.8)
    add_text(
        slide,
        f"{result:.2f}",
        x + 0.25,
        4.95,
        1.6,
        0.48,
        size=25,
        color=accent,
        bold=True,
        margin=0,
    )
    add_text(
        slide,
        "MAE (W/m²)",
        x + 0.25,
        5.45,
        1.5,
        0.2,
        size=8.5,
        color=MUTED,
        margin=0,
    )
    add_text(
        slide,
        f"{rank}º geral",
        x + 2.08,
        5.13,
        1.4,
        0.24,
        size=11,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.RIGHT,
        margin=0,
    )


def slide_baselines(prs, metrics):
    slide = blank_slide(prs)
    add_header(
        slide,
        "Referências simples",
        "Três baselines: o teste mínimo que todo modelo precisa superar",
        "Elas medem quanto valor real a complexidade acrescenta ao padrão sazonal.",
        8,
        TEAL,
    )
    baseline_card(
        slide,
        0.7,
        "Persistência",
        "Repete a GHI observada no mês imediatamente anterior.",
        metrics["Persistencia"]["mae"],
        10,
        ORANGE,
        "t−1",
    )
    baseline_card(
        slide,
        4.8,
        "Sazonal ingênuo",
        "Usa a observação do mesmo mês do ano anterior.",
        metrics["Sazonal ingenuo"]["mae"],
        8,
        BLUE,
        "t−12",
    )
    baseline_card(
        slide,
        8.9,
        "Climatologia",
        "Média histórica do mês e da localidade, calculada apenas no treino.",
        metrics["Climatologia"]["mae"],
        1,
        LIME,
        "μm",
    )
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        0.7,
        6.14,
        11.94,
        0.68,
        fill=CARD_ALT,
        line=LIME,
    )
    add_text(
        slide,
        "Resultado-chave: a climatologia foi o melhor método geral — a sazonalidade anual dominou este conjunto mensal.",
        0.97,
        6.34,
        11.4,
        0.26,
        size=13.2,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    add_footer(slide, "Baselines")


def slide_ranking(prs, metrics):
    slide = blank_slide(prs)
    add_header(
        slide,
        "Resultados pontuais",
        "Ranking por macro-MAE: simples e recorrentes lideram",
        "Média não ponderada do MAE nas dez localidades; menor é melhor.",
        9,
        TEAL,
    )
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        0.7,
        1.87,
        8.2,
        4.98,
        fill="FFFFFF",
        line="FFFFFF",
    )
    slide.shapes.add_picture(
        str(RESULTADOS / "figuras" / "mae_medio_modelos.png"),
        Inches(0.9),
        Inches(2.03),
        width=Inches(7.8),
    )
    add_info_card(
        slide,
        9.18,
        1.87,
        3.45,
        1.38,
        "1º lugar",
        "Climatologia\n12,07 W/m²",
        LIME,
        body_size=15,
    )
    add_info_card(
        slide,
        9.18,
        3.46,
        3.45,
        1.38,
        "Melhor modelo aprendido",
        "LSTM\n12,22 W/m²",
        TEAL,
        body_size=15,
    )
    add_info_card(
        slide,
        9.18,
        5.05,
        3.45,
        1.8,
        "Leitura",
        "LSTM ficou apenas 0,15 W/m² atrás da climatologia. RNN e DilatedRNN completaram o grupo líder.",
        PURPLE,
        body_size=12.4,
    )
    add_footer(slide, "Ranking")


def add_prob_row(slide, y, name, crps, coverage, width, accent):
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        0.92,
        y,
        7.15,
        0.83,
        fill=CARD_ALT,
        line=GRID,
    )
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0.92,
        y,
        0.08,
        0.83,
        fill=accent,
        line=accent,
    )
    add_text(slide, name, 1.17, y + 0.25, 1.25, 0.24, size=13, bold=True, color=accent, margin=0)
    for x, value, suffix in (
        (2.75, f"{crps:.2f}", "CRPS"),
        (4.45, f"{coverage:.2f}%", "cobertura"),
        (6.22, f"{width:.2f}", "largura"),
    ):
        add_text(slide, value, x, y + 0.12, 1.35, 0.28, size=15, bold=True, color=WHITE, margin=0)
        add_text(slide, suffix, x, y + 0.48, 1.35, 0.16, size=7.5, color=MUTED, margin=0)


def slide_interpretacao(prs, prob_metrics):
    slide = blank_slide(prs)
    add_header(
        slide,
        "Resultados e conclusão",
        "Precisão, cobertura e largura contam histórias diferentes",
        "A avaliação probabilística exige olhar calibração e nitidez em conjunto.",
        10,
        PINK,
    )
    add_text(
        slide,
        "RESULTADOS PROBABILÍSTICOS",
        0.92,
        1.95,
        3.3,
        0.22,
        size=9,
        color=PINK,
        bold=True,
        margin=0,
    )
    d = prob_metrics["DeepAR"]
    n = prob_metrics["DeepNPTS"]
    add_prob_row(slide, 2.28, "DeepAR", d["crps"], d["picp"], d["mpiw"], PINK)
    add_prob_row(slide, 3.33, "DeepNPTS", n["crps"], n["picp"], n["mpiw"], YELLOW)
    add_text(
        slide,
        "CRPS e largura em W/m² · intervalo nominal de 90%",
        0.93,
        4.31,
        4.0,
        0.18,
        size=7.5,
        color=MUTED,
        margin=0,
    )
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        8.39,
        1.87,
        4.24,
        3.01,
        fill=CARD,
        line=GRID,
    )
    add_text(
        slide,
        "VENCEDORES POR LOCALIDADE",
        8.66,
        2.13,
        3.7,
        0.22,
        size=9,
        color=TEAL,
        bold=True,
        margin=0,
    )
    winners = [
        ("LSTM", 4, LIME),
        ("Climatologia", 3, TEAL),
        ("RNN", 1, BLUE),
        ("DeepAR", 1, PINK),
        ("Sazonal", 1, ORANGE),
    ]
    for idx, (name, count, color) in enumerate(winners):
        yy = 2.65 + idx * 0.4
        add_text(slide, name, 8.66, yy, 1.25, 0.18, size=9.5, color=WHITE, margin=0)
        add_shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            9.93,
            yy + 0.01,
            count * 0.43,
            0.16,
            fill=color,
            line=color,
        )
        add_text(
            slide,
            str(count),
            11.92,
            yy,
            0.25,
            0.18,
            size=9.5,
            color=color,
            bold=True,
            align=PP_ALIGN.RIGHT,
            margin=0,
        )
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        0.7,
        5.19,
        11.93,
        1.64,
        fill=CARD_ALT,
        line=GRID,
    )
    takeaways = [
        ("01", "Sazonalidade forte", "A climatologia venceu o macro-MAE.", LIME),
        ("02", "Recorrência competitiva", "LSTM venceu 4 das 10 localidades.", TEAL),
        ("03", "Incerteza tem trade-off", "DeepAR foi mais nítido; DeepNPTS cobriu mais.", PINK),
    ]
    for idx, (number, title, body, color) in enumerate(takeaways):
        x = 0.98 + idx * 3.88
        add_circle(slide, x, 5.51, 0.32, fill=color)
        add_text(
            slide,
            number,
            x,
            5.59,
            0.32,
            0.14,
            size=7.5,
            color=BG,
            bold=True,
            align=PP_ALIGN.CENTER,
            margin=0,
        )
        add_text(slide, title, x + 0.48, 5.47, 2.95, 0.24, size=11.5, color=WHITE, bold=True, margin=0)
        add_text(slide, body, x + 0.48, 5.86, 2.95, 0.48, size=10, color=MUTED, margin=0)
    add_footer(slide, "Conclusões")


def load_metrics():
    metrics = {}
    with (RESULTADOS / "metricas_medias_modelos.csv").open(
        encoding="utf-8", newline=""
    ) as handle:
        for row in csv.DictReader(handle):
            metrics[row["Modelo"]] = {
                "mae": float(row["MAE_media_wm2"]),
                "rmse": float(row["RMSE_media_wm2"]),
            }
    prob = {}
    with (RESULTADOS / "metricas_probabilisticas_medias.csv").open(
        encoding="utf-8", newline=""
    ) as handle:
        for row in csv.DictReader(handle):
            prob[row["Modelo"]] = {
                "crps": float(row["CRPS_medio_wm2"]),
                "picp": float(row["PICP_medio_percentual"]),
                "mpiw": float(row["MPIW_medio_wm2"]),
            }
    return metrics, prob


def main():
    metrics, prob_metrics = load_metrics()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.core_properties.title = "Modelos, baselines e resultados — previsão mensal de GHI"
    prs.core_properties.subject = "Síntese do experimento canônico do TCC"
    prs.core_properties.author = "Projeto TCC"
    prs.core_properties.keywords = "GHI, machine learning, baselines, DeepNPTS, DeepAR"

    slide_xgboost(prs, metrics)
    slide_mlp(prs, metrics)
    slide_rnn(prs, metrics)
    slide_lstm(prs, metrics)
    slide_dilated(prs, metrics)
    slide_deepar(prs, metrics, prob_metrics)
    slide_deepnpts(prs, metrics, prob_metrics)
    slide_baselines(prs, metrics)
    slide_ranking(prs, metrics)
    slide_interpretacao(prs, prob_metrics)

    SAIDA.parent.mkdir(parents=True, exist_ok=True)
    prs.save(SAIDA)
    print(SAIDA)


if __name__ == "__main__":
    main()
