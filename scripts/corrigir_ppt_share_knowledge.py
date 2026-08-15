#!/usr/bin/env python3
"""Corrige cirurgicamente o deck Share Knowledge preservando o visual original."""

from __future__ import annotations

from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "powerpoint" / "sidi_automation_ml_forecasting_pt.pptx"
OUTPUT = ROOT / "powerpoint" / "sidi_automation_ml_forecasting_pt_corrigido.pptx"

W = 13.333333
H = 7.5
WHITE = "FFFFFF"
BLACK = "111111"
DARK = "272727"
GRAY = "68706C"
MID = "A2A8A5"
GREEN = "008B3D"
DARK_GREEN = "007334"
LIGHT_GREEN = "EAF6ED"
LIGHT_GRAY = "F4F5F4"
ORANGE = "E88200"
RED = "D71920"
BLUE = "2775B6"
FONT = "Arial"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def rect(slide, x, y, w, h, *, fill=WHITE, line=None, rounded=False, radius=None):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1.0)
    else:
        shape.line.fill.background()
    return shape


def text_box(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    size=12,
    color=BLACK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
    rotation=0,
    margin=0.02,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    shape.rotation = rotation
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.line_spacing = 1.0
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return shape


def cover_and_text(slide, text, x, y, w, h, **kwargs):
    rect(slide, x, y, w, h, fill=kwargs.pop("fill", WHITE))
    return text_box(slide, text, x, y, w, h, **kwargs)


def pill(slide, text, x, y, w, h=0.27, *, fill=LIGHT_GREEN, color=DARK_GREEN, size=8.5):
    rect(slide, x, y, w, h, fill=fill, line=GREEN, rounded=True)
    text_box(
        slide,
        text,
        x,
        y,
        w,
        h,
        size=size,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def set_notes(slide, content: str) -> None:
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    tf.paragraphs[0].text = content.strip()


def add_line(slide, x1, y1, x2, y2, *, color=MID, width=1.0):
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)
    return shape


def add_arrow(slide, x1, y1, x2, y2, *, color=GREEN, width=1.8):
    shape = add_line(slide, x1, y1, x2, y2, color=color, width=width)
    try:
        shape.line.end_arrowhead = True
    except (AttributeError, ValueError):
        pass
    return shape


def rich_text(slide, parts, x, y, w, h, *, size=22, align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    for value, color, bold in parts:
        run = p.add_run()
        run.text = value
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return shape


def add_marker(slide, x, y, *, color, size=0.09):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x - size / 2),
        Inches(y - size / 2),
        Inches(size),
        Inches(size),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.color.rgb = rgb(WHITE)
    shape.line.width = Pt(0.6)


def card(slide, x, y, w, h, title, body, *, accent=GREEN):
    rect(slide, x, y, w, h, fill=WHITE, line="DDE6DF", rounded=True)
    rect(slide, x, y, 0.06, h, fill=accent, rounded=True)
    text_box(slide, title, x + 0.20, y + 0.10, w - 0.30, 0.30, size=10.5, color=accent, bold=True)
    text_box(slide, body, x + 0.20, y + 0.43, w - 0.30, h - 0.53, size=10.5, color=DARK, valign=MSO_ANCHOR.TOP)


def correct_existing_slides(prs: Presentation) -> None:
    # 1 — grafia institucional na capa.
    slide = prs.slides[0]
    cover_and_text(
        slide,
        "SiDi Automation",
        0.58,
        1.64,
        5.30,
        0.76,
        size=35,
        color=BLACK,
        valign=MSO_ANCHOR.MIDDLE,
    )

    # 2 — o desenho é conceitual; remove a unidade de energia incorreta.
    slide = prs.slides[1]
    cover_and_text(
        slide,
        "Índice ilustrativo de GHI",
        0.45,
        3.01,
        0.44,
        1.43,
        size=8.5,
        color=DARK,
        bold=True,
        align=PP_ALIGN.CENTER,
        rotation=270,
    )
    cover_and_text(
        slide,
        "Índice ilustrativo de GHI",
        8.39,
        3.03,
        0.44,
        1.42,
        size=8.5,
        color=DARK,
        bold=True,
        align=PP_ALIGN.CENTER,
        rotation=270,
    )

    # 3 — linguagem e status do TimesNet.
    slide = prs.slides[2]
    cover_and_text(
        slide,
        "O tempo representado por\natributos derivados",
        0.63,
        5.07,
        2.43,
        0.43,
        size=10.5,
        color=DARK,
        align=PP_ALIGN.CENTER,
    )
    cover_and_text(
        slide,
        "PREVISÃO HORÁRIA",
        9.68,
        2.64,
        1.70,
        0.28,
        size=10.5,
        color=DARK_GREEN,
        bold=True,
        align=PP_ALIGN.LEFT,
    )
    cover_and_text(
        slide,
        "Modelo horário de previsão\ndireta multietapas",
        9.46,
        5.08,
        2.20,
        0.43,
        size=10.2,
        color=DARK,
        align=PP_ALIGN.CENTER,
    )
    cover_and_text(
        slide,
        "Persistência   ·   Sazonal ingênuo   ·   Climatologia",
        4.86,
        5.83,
        3.95,
        0.25,
        size=10.2,
        color=DARK,
        align=PP_ALIGN.CENTER,
    )

    # 4 — redação natural da limitação tabular.
    slide = prs.slides[3]
    cover_and_text(
        slide,
        "Atenção:  Depende de boa engenharia de atributos temporais",
        0.99,
        5.37,
        4.78,
        0.28,
        size=9.6,
        color=ORANGE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 9 — explicita que a figura é didática e que o contexto real é mensal.
    slide = prs.slides[8]
    cover_and_text(
        slide,
        "No experimento, os candidatos são as 12 posições do contexto mensal.",
        0.75,
        5.83,
        9.05,
        0.24,
        size=8.2,
        color=DARK_GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 10 — status final do experimento horário.
    slide = prs.slides[9]
    cover_and_text(
        slide,
        "De Sinais 1D a Padrões Temporais 2D: TimesNet",
        0.39,
        1.39,
        10.18,
        0.48,
        size=25.0,
        color=BLACK,
        bold=True,
    )
    rect(slide, 9.52, 5.65, 3.35, 0.80, fill=WHITE, line=GREEN, rounded=True)
    text_box(
        slide,
        "AVALIAÇÃO CONCLUÍDA",
        9.72,
        5.78,
        2.95,
        0.24,
        size=11.5,
        color=DARK_GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    text_box(
        slide,
        "Resultados apresentados separadamente\ndo ranking mensal",
        9.72,
        6.02,
        2.95,
        0.32,
        size=8.8,
        color=DARK,
        align=PP_ALIGN.CENTER,
    )

    # 11 — definição exata do nRMSE e leitura probabilística.
    slide = prs.slides[10]
    cover_and_text(
        slide,
        "RMSE dividido pela média observada\nda localidade e apresentado em %.\nMenor é melhor.",
        1.97,
        4.91,
        2.00,
        0.86,
        size=9.2,
        color=DARK,
        valign=MSO_ANCHOR.TOP,
    )
    cover_and_text(
        slide,
        "CRPS: menor é melhor. Cobertura:\npróxima do nível nominal. Largura:\nmenor quando a cobertura é comparável.",
        9.36,
        4.91,
        2.02,
        0.86,
        size=8.8,
        color=DARK,
        valign=MSO_ANCHOR.TOP,
    )

    # 12 — termos, cautela estatística e status do TimesNet.
    slide = prs.slides[11]
    cover_and_text(
        slide,
        "Sazonal\ningênuo",
        0.53,
        4.45,
        1.08,
        0.34,
        size=8.4,
        color=BLACK,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )
    cover_and_text(
        slide,
        "MENOR MACRO-MAE",
        8.42,
        2.35,
        2.05,
        0.28,
        size=10.5,
        color=DARK_GREEN,
        bold=True,
    )
    cover_and_text(
        slide,
        "RESULTADO DESCRITIVO",
        8.42,
        4.74,
        2.18,
        0.28,
        size=10.3,
        color=DARK_GREEN,
        bold=True,
    )
    cover_and_text(
        slide,
        "Recorrentes ficaram entre os quatro\nprimeiros, mas sem diferença\nconclusiva no topo.",
        8.42,
        5.08,
        2.54,
        0.54,
        size=8.8,
        color=DARK,
        valign=MSO_ANCHOR.TOP,
    )
    rect(slide, 0.41, 6.05, 10.82, 0.57, fill=WHITE, line=GREEN, rounded=True)
    text_box(
        slide,
        "LSTM × climatologia: Δ = +0,15 W/m²; IC95% [−0,95; 1,32]; diferença inconclusiva.",
        1.00,
        6.10,
        9.72,
        0.20,
        size=8.2,
        color=DARK_GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    text_box(
        slide,
        "TimesNet é um experimento horário separado; os resultados não são diretamente comparáveis.",
        1.00,
        6.34,
        9.72,
        0.18,
        size=7.8,
        color=DARK,
        align=PP_ALIGN.CENTER,
    )

    # 13 original — conclusões e próximo passo atualizados; a numeração será corrigida após inserir o novo slide.
    slide = prs.slides[12]
    cover_and_text(
        slide,
        "• Baselines fortes são essenciais.\n• A forma de representar o tempo influenciou o ranking observado.\n• Complexidade não garante melhor desempenho.\n• Cobertura deve ser avaliada junto com a largura do intervalo.\n• Não existe um modelo universalmente vencedor.",
        1.30,
        5.05,
        4.38,
        1.10,
        size=9.0,
        color=DARK,
        valign=MSO_ANCHOR.TOP,
    )
    cover_and_text(
        slide,
        "Repetir o experimento horário com múltiplas sementes e investigar variáveis meteorológicas adicionais.",
        7.35,
        5.35,
        3.62,
        0.73,
        size=9.4,
        color=DARK,
        valign=MSO_ANCHOR.TOP,
    )


def add_hourly_results_slide(prs: Presentation):
    # Usa a imagem do slide mensal como base para preservar logo, rodapé e grafismo SiDi.
    background_blob = prs.slides[11].shapes[0].image.blob
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.add_picture(BytesIO(background_blob), 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Limpa apenas a área de conteúdo, mantendo a identidade visual externa.
    rect(slide, 0.34, 1.04, 10.95, 5.65, fill=WHITE)
    text_box(slide, "RESULTADOS HORÁRIOS", 0.42, 1.13, 3.30, 0.25, size=10.5, color=DARK_GREEN, bold=True)
    text_box(
        slide,
        "XGBoost liderou e TimesNet foi competitivo",
        0.42,
        1.42,
        10.65,
        0.47,
        size=25.5,
        color=BLACK,
        bold=True,
    )
    text_box(slide, "12", 12.78, 0.22, 0.25, 0.20, size=10.5, color=DARK_GREEN, bold=True, align=PP_ALIGN.RIGHT)

    # Área do gráfico.
    rect(slide, 0.43, 2.06, 7.38, 4.20, fill=WHITE, line="DDE6DF", rounded=True)
    text_box(
        slide,
        "RMSE no teste de 2024 — todas as horas, pós-processado — menor é melhor",
        0.74,
        2.21,
        6.82,
        0.28,
        size=10.8,
        color=DARK,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    plot_left, plot_right = 1.12, 7.28
    plot_top, plot_bottom = 3.04, 5.62
    y_min, y_max = 90.0, 150.0

    def y_pos(value):
        return plot_bottom - (value - y_min) / (y_max - y_min) * (plot_bottom - plot_top)

    # Grade e eixos.
    for tick in (90, 100, 110, 120, 130, 140, 150):
        y = y_pos(tick)
        add_line(slide, plot_left, y, plot_right, y, color="D9DEDB", width=0.7)
        text_box(slide, str(tick), 0.72, y - 0.10, 0.31, 0.20, size=8.2, color=GRAY, align=PP_ALIGN.RIGHT)
    add_line(slide, plot_left, plot_top, plot_left, plot_bottom, color=GRAY, width=1.0)
    add_line(slide, plot_left, plot_bottom, plot_right, plot_bottom, color=GRAY, width=1.0)
    text_box(slide, "RMSE (W/m²)", 0.47, 3.73, 0.40, 1.10, size=8.4, color=GRAY, bold=True, rotation=270, align=PP_ALIGN.CENTER)

    xs = [2.00, 4.18, 6.36]
    for x, label in zip(xs, ("24 h", "48 h", "72 h"), strict=True):
        text_box(slide, label, x - 0.33, 5.72, 0.66, 0.24, size=9.2, color=DARK, bold=True, align=PP_ALIGN.CENTER)

    series = [
        ("XGBoost", [96.05, 98.23, 99.22], DARK_GREEN),
        ("TimesNet", [97.21, 100.26, 101.04], "39B54A"),
        ("LSTM", [100.01, 103.00, 104.68], BLUE),
        ("Persistência", [122.21, 128.35, 131.32], ORANGE),
        ("Sazonal ingênuo", [142.78, 142.79, 142.76], "7B827E"),
    ]
    legend_x = [0.80, 2.05, 3.22, 4.24, 5.62]
    for (name, values, color), lx in zip(series, legend_x, strict=True):
        add_line(slide, lx, 2.73, lx + 0.25, 2.73, color=color, width=2.2)
        text_box(slide, name, lx + 0.30, 2.61, 1.15, 0.24, size=7.9, color=DARK, bold=name in ("XGBoost", "TimesNet"))
        points = [(x, y_pos(value)) for x, value in zip(xs, values, strict=True)]
        for (x1, y1), (x2, y2) in zip(points, points[1:]):
            add_line(slide, x1, y1, x2, y2, color=color, width=2.0 if name in ("XGBoost", "TimesNet") else 1.35)
        for x, y in points:
            add_marker(slide, x, y, color=color, size=0.11 if name in ("XGBoost", "TimesNet") else 0.08)

    # Leituras principais à direita.
    card(
        slide,
        8.05,
        2.13,
        3.05,
        1.04,
        "MENOR RMSE NOS 3 HORIZONTES",
        "XGBoost: 96,05 · 98,23 · 99,22 W/m²",
        accent=DARK_GREEN,
    )
    card(
        slide,
        8.05,
        3.37,
        3.05,
        1.04,
        "TIMESNET COMPETITIVO",
        "Segundo menor RMSE em 24 h, 48 h e 72 h.",
        accent="39B54A",
    )
    card(
        slide,
        8.05,
        4.61,
        3.05,
        1.17,
        "LEITURA CORRETA",
        "Experimento horário separado e executado com uma única semente: 42.",
        accent=ORANGE,
    )
    rect(slide, 0.67, 6.04, 10.16, 0.40, fill=LIGHT_GREEN, line="C8E4CF", rounded=True)
    text_box(
        slide,
        "Os erros horários não são diretamente comparáveis ao ranking mensal.",
        0.88,
        6.08,
        9.72,
        0.30,
        size=10.2,
        color=DARK_GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    set_notes(
        slide,
        """
Apresentar esta lâmina como um experimento separado do ranking mensal.
Os resultados horários já foram concluídos e foram calculados no teste de 2024, considerando todas as horas e as previsões pós-processadas.
O XGBoost apresentou o menor RMSE nos horizontes de 24, 48 e 72 horas. O TimesNet ficou próximo e ocupou a segunda posição nos três horizontes, portanto deve ser descrito como competitivo, não como vencedor.
Destacar a limitação: o experimento horário utilizou apenas a semente 42, então ainda falta medir a variabilidade entre inicializações.
Transição: depois da comparação pontual, voltar aos modelos probabilísticos e discutir calibração e largura de intervalos.
""",
    )

    # Insere antes do slide final original.
    slide_ids = prs.slides._sldIdLst
    new_id = slide_ids[-1]
    slide_ids.remove(new_id)
    slide_ids.insert(len(slide_ids) - 1, new_id)
    return slide


def add_intro_background(prs: Presentation):
    background_blob = prs.slides[11].shapes[0].image.blob
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.add_picture(BytesIO(background_blob), 0, 0, width=prs.slide_width, height=prs.slide_height)
    rect(slide, 0.32, 1.02, 10.97, 5.70, fill=WHITE)
    return slide


def add_need_slide(prs: Presentation):
    slide = add_intro_background(prs)
    text_box(slide, "NECESSIDADE DO PROJETO", 0.44, 1.16, 3.5, 0.23, size=10.5, color=DARK_GREEN, bold=True)
    text_box(
        slide,
        "Por que prever a disponibilidade do recurso solar?",
        0.44,
        1.40,
        10.55,
        0.48,
        size=24.5,
        color=BLACK,
        bold=True,
    )

    xs = [0.45, 3.43, 6.41, 9.39]
    titles = [
        "Variabilidade da irradiância",
        "Incerteza sobre\no recurso solar",
        "Previsão",
        "Apoio ao planejamento\ne à análise",
    ]
    for x, title in zip(xs, titles, strict=True):
        rect(slide, x, 2.05, 2.55, 2.65, fill=WHITE, line="C9D4CD", rounded=True)
        text_box(slide, title, x + 0.18, 2.17, 2.19, 0.45, size=11.2, color=DARK_GREEN, bold=True, align=PP_ALIGN.CENTER)

    # Cartão 1: variação temporal e sazonal.
    x = xs[0]
    for cx, fill in ((x + 0.50, "F9CF42"), (x + 1.28, WHITE), (x + 2.05, "D5D8D6")):
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - 0.16), Inches(2.79), Inches(0.32), Inches(0.32))
        circ.fill.solid(); circ.fill.fore_color.rgb = rgb(fill); circ.line.color.rgb = rgb(GREEN if fill == WHITE else MID)
    add_arrow(slide, x + 0.75, 2.95, x + 1.02, 2.95, color=DARK_GREEN, width=1.2)
    add_arrow(slide, x + 1.53, 2.95, x + 1.80, 2.95, color=DARK_GREEN, width=1.2)
    points = [(x + 0.25, 3.78), (x + 0.63, 3.38), (x + 1.02, 3.68), (x + 1.41, 3.27), (x + 1.80, 3.70), (x + 2.23, 3.42)]
    for a, b in zip(points, points[1:]): add_line(slide, *a, *b, color=GREEN, width=1.6)
    text_box(slide, "Primavera        Verão        Outono       Inverno", x + 0.18, 4.15, 2.19, 0.23, size=7.3, color=GRAY, align=PP_ALIGN.CENTER)

    # Cartão 2: incerteza.
    x = xs[1]
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.90), Inches(2.91), Inches(0.75), Inches(0.75))
    circ.fill.solid(); circ.fill.fore_color.rgb = rgb(WHITE); circ.line.color.rgb = rgb(GREEN); circ.line.width = Pt(2)
    text_box(slide, "?", x + 0.90, 2.91, 0.75, 0.75, size=28, color=DARK_GREEN, bold=True, align=PP_ALIGN.CENTER)
    for dx, dy in ((0.40, 2.78), (1.90, 2.78), (0.40, 3.75), (1.90, 3.75)):
        text_box(slide, "?", x + dx, dy, 0.25, 0.25, size=13, color=MID, bold=True, align=PP_ALIGN.CENTER)
        add_line(slide, x + 1.27, 3.28, x + dx + 0.12, dy + 0.12, color=MID, width=0.9)

    # Cartão 3: rede e faixa de previsão.
    x = xs[2]
    nodes_left = [(x + 0.36, 2.88), (x + 0.36, 3.25), (x + 0.36, 3.62)]
    nodes_right = [(x + 1.00, 2.95), (x + 1.00, 3.25), (x + 1.00, 3.55)]
    for a in nodes_left:
        for b in nodes_right: add_line(slide, *a, *b, color=MID, width=0.7)
    for nx, ny in nodes_left + nodes_right:
        node = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(nx - 0.06), Inches(ny - 0.06), Inches(0.12), Inches(0.12))
        node.fill.solid(); node.fill.fore_color.rgb = rgb(WHITE); node.line.color.rgb = rgb(GREEN)
    forecast = [(x + 1.35, 3.52), (x + 1.62, 3.34), (x + 1.90, 3.42), (x + 2.18, 3.18)]
    for a, b in zip(forecast, forecast[1:]): add_line(slide, *a, *b, color=DARK_GREEN, width=1.4)
    for offset, shade in ((0.18, "D9F0DF"), (0.09, "BDE5C7")):
        for a, b in zip(forecast, forecast[1:]):
            add_line(slide, a[0], a[1] - offset, b[0], b[1] - offset, color=shade, width=4.0)
            add_line(slide, a[0], a[1] + offset, b[0], b[1] + offset, color=shade, width=4.0)
    text_box(slide, "Dados históricos      Modelos de ML      Previsões", x + 0.18, 4.12, 2.19, 0.28, size=7.5, color=DARK, align=PP_ALIGN.CENTER)

    # Cartão 4: decisões e análises.
    x = xs[3]
    for ox, oy in ((0.32, 2.78), (1.34, 2.78), (0.32, 3.48), (1.34, 3.48)):
        rect(slide, x + ox, oy, 0.82, 0.55, fill=LIGHT_GRAY, line="D9DEDB", rounded=True)
    for i, h in enumerate((0.15, 0.28, 0.39)):
        rect(slide, x + 1.55 + i * 0.18, 3.16 - h, 0.10, h, fill="6AC07B")
    for i in range(3):
        rect(slide, x + 0.48, 3.60 + i * 0.13, 0.09, 0.09, fill=WHITE, line=GREEN)
        add_line(slide, x + 0.65, 3.65 + i * 0.13, x + 0.96, 3.65 + i * 0.13, color=MID, width=0.8)
    text_box(slide, "Decisões mais informadas\npara estudos e estratégias", x + 0.28, 4.12, 1.99, 0.34, size=8.4, color=DARK, align=PP_ALIGN.CENTER)

    # Setas entre os cartões.
    for x in (3.10, 6.08, 9.06): add_arrow(slide, x, 3.33, x + 0.22, 3.33, color=GREEN, width=2.0)

    explanations = [
        "A irradiância varia ao longo do\ntempo e entre localidades",
        "Médias históricas ajudam, mas podem\nnão refletir mudanças recentes",
        "Modelos transformam o histórico\nem previsões de valores futuros",
        "O projeto compara modelos aprendidos\ncom referências simples",
    ]
    for x, body in zip(xs, explanations, strict=True):
        text_box(slide, body, x + 0.05, 4.83, 2.45, 0.45, size=9.2, color=DARK, align=PP_ALIGN.CENTER)

    rect(slide, 3.30, 5.62, 6.85, 0.73, fill=WHITE, line=GREEN, rounded=True)
    text_box(slide, "ESCOPO DO ESTUDO", 3.55, 5.73, 1.75, 0.20, size=9.6, color=DARK_GREEN, bold=True)
    text_box(
        slide,
        "O estudo prevê o recurso solar nas coordenadas analisadas. Não prevê consumo industrial nem geração fotovoltaica. As fábricas são apenas referências geográficas.",
        3.55,
        5.96,
        6.32,
        0.29,
        size=8.4,
        color=DARK,
    )
    set_notes(slide, """
Usar esta lâmina apenas para contextualizar o caso prático do Share Knowledge.
A irradiância varia no tempo e entre localidades; por isso, previsões podem ser mais úteis do que depender somente de uma média histórica.
O objetivo do projeto é comparar formas diferentes de representar o tempo e verificar se modelos aprendidos superam referências simples.
Ressalvar que as fábricas são apenas referências geográficas e que não existem dados internos de consumo ou geração.
Transição: definir rapidamente qual variável está sendo prevista.
""")
    return slide


def add_ghi_slide(prs: Presentation):
    slide = add_intro_background(prs)
    text_box(slide, "O QUE É GHI?", 0.45, 1.50, 2.4, 0.23, size=10.5, color=DARK_GREEN, bold=True)
    text_box(slide, "O que estamos prevendo?", 0.45, 1.10, 7.0, 0.48, size=27, color=BLACK)

    # Sol e componentes da irradiância.
    sun = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), Inches(2.52), Inches(0.48), Inches(0.48))
    sun.fill.solid(); sun.fill.fore_color.rgb = rgb("F9B000"); sun.line.fill.background()
    for dx, dy in ((0.0, -0.20), (0.0, 0.68), (-0.20, 0.24), (0.68, 0.24), (-0.14, -0.10), (0.58, -0.10), (-0.14, 0.58), (0.58, 0.58)):
        add_line(slide, 1.09, 2.76, 0.85 + dx + 0.24, 2.52 + dy + 0.24, color="F9B000", width=1.4)
    add_line(slide, 3.10, 2.48, 3.10, 5.42, color=GRAY, width=1.0)
    text_box(slide, "Zênite", 2.83, 2.30, 0.55, 0.22, size=8.5, color=GRAY, align=PP_ALIGN.CENTER)
    add_arrow(slide, 1.35, 3.05, 3.10, 5.35, color="F0A000", width=2.1)
    text_box(slide, "DNI", 1.12, 3.55, 0.55, 0.25, size=14, color="F0A000", bold=True, align=PP_ALIGN.CENTER)
    text_box(slide, "Irradiância\nDireta Normal", 0.84, 3.82, 1.08, 0.42, size=8.2, color=DARK, align=PP_ALIGN.CENTER)
    for sx, sy in ((3.60, 3.05), (4.45, 3.15), (5.30, 3.38), (3.45, 4.18), (4.85, 4.42)):
        add_arrow(slide, sx, sy, sx - 0.25, sy + 0.80, color=BLUE, width=1.2)
    text_box(slide, "DHI", 4.32, 2.68, 0.60, 0.25, size=14, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    text_box(slide, "Irradiância\nDifusa Horizontal", 4.03, 2.94, 1.18, 0.40, size=8.2, color=DARK, align=PP_ALIGN.CENTER)
    add_line(slide, 0.95, 5.42, 5.70, 5.42, color=GRAY, width=6.0)
    text_box(slide, "Superfície horizontal", 2.25, 5.55, 2.10, 0.24, size=10.5, color=DARK, align=PP_ALIGN.CENTER)
    text_box(slide, "θz", 2.60, 3.32, 0.45, 0.28, size=15, color=DARK_GREEN, bold=True, align=PP_ALIGN.CENTER)

    # Fórmula e definição.
    rect(slide, 6.45, 2.08, 4.23, 0.82, fill=WHITE, line=GREEN, rounded=True)
    rich_text(
        slide,
        [("GHI", DARK_GREEN, True), (" = ", BLACK, False), ("DHI", BLUE, True), (" + ", BLACK, False), ("DNI", "F0A000", True), (" × cos(θz)", BLACK, False)],
        6.67,
        2.25,
        3.80,
        0.42,
        size=19,
    )
    rich_text(slide, [("GHI", DARK_GREEN, True), (" é a potência solar total que chega a uma superfície horizontal.", BLACK, False)], 6.50, 3.05, 4.20, 0.55, size=12, align=PP_ALIGN.LEFT)
    rich_text(slide, [("Unidade: ", GRAY, False), ("W/m²", DARK_GREEN, True)], 6.50, 3.58, 2.4, 0.30, size=12, align=PP_ALIGN.LEFT)

    # Comparação irradiância versus energia acumulada.
    rect(slide, 6.45, 4.05, 2.35, 2.18, fill=WHITE, line=GREEN, rounded=True)
    rect(slide, 9.05, 4.05, 2.50, 2.18, fill=WHITE, line=MID, rounded=True)
    meter = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.37), Inches(4.30), Inches(0.52), Inches(0.52))
    meter.fill.solid(); meter.fill.fore_color.rgb = rgb(WHITE); meter.line.color.rgb = rgb(GREEN); meter.line.width = Pt(2)
    add_line(slide, 7.63, 4.56, 7.77, 4.43, color=GREEN, width=1.6)
    text_box(slide, "IRRADIÂNCIA", 6.78, 4.92, 1.70, 0.28, size=11, color=DARK_GREEN, bold=True, align=PP_ALIGN.CENTER)
    text_box(slide, "potência solar\ninstantânea ou média", 6.75, 5.20, 1.75, 0.40, size=9.4, color=DARK, align=PP_ALIGN.CENTER)
    text_box(slide, "unidade: W/m²", 6.82, 5.63, 1.60, 0.24, size=9.0, color=DARK_GREEN, bold=True, align=PP_ALIGN.CENTER)
    rect(slide, 6.45, 5.90, 2.35, 0.33, fill=LIGHT_GREEN)
    text_box(slide, "✓  alvo deste projeto", 6.67, 5.94, 1.90, 0.24, size=9.2, color=DARK_GREEN, bold=True, align=PP_ALIGN.CENTER)

    text_box(slide, "IRRADIAÇÃO OU\nENERGIA ACUMULADA", 9.32, 4.75, 1.96, 0.48, size=10.2, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    text_box(slide, "energia recebida ao\nlongo de um período", 9.35, 5.20, 1.90, 0.40, size=9.4, color=DARK, align=PP_ALIGN.CENTER)
    text_box(slide, "unidade: Wh/m² ou kWh/m²", 9.18, 5.62, 2.24, 0.24, size=8.4, color=GRAY, align=PP_ALIGN.CENTER)
    rect(slide, 9.05, 5.90, 2.50, 0.33, fill=LIGHT_GRAY)
    text_box(slide, "✕  não é o alvo deste projeto", 9.20, 5.94, 2.20, 0.24, size=8.8, color=GRAY, bold=True, align=PP_ALIGN.CENTER)

    set_notes(slide, """
Definir GHI rapidamente, sem transformar a apresentação em uma aula de energia solar.
GHI é a potência solar total que chega a uma superfície horizontal e combina a componente direta projetada na horizontal com a componente difusa.
O alvo do projeto é irradiância média em W/m². Não é energia acumulada, que seria expressa em Wh/m² ou kWh/m².
Transição: mostrar como os experimentos mensal e horário foram avaliados.
""")
    return slide


def add_protocol_slide(prs: Presentation):
    slide = add_intro_background(prs)
    text_box(slide, "DADOS E PROTOCOLO EXPERIMENTAL", 0.43, 1.55, 4.3, 0.23, size=10.5, color=DARK_GREEN, bold=True)
    text_box(slide, "Como a previsão foi avaliada?", 0.43, 1.08, 8.0, 0.50, size=27, color=BLACK, bold=True)

    # Experimento mensal.
    rect(slide, 0.42, 2.00, 8.45, 4.52, fill=WHITE, line=GREEN, rounded=True)
    text_box(slide, "EXPERIMENTO MENSAL", 0.80, 2.16, 2.65, 0.28, size=14, color=BLACK, bold=True)
    text_box(slide, "NSRDB/NREL — GOES Aggregated PSM v4", 0.80, 2.45, 3.65, 0.22, size=9.2, color=DARK)

    # Bloco de localidades.
    rect(slide, 0.62, 2.78, 2.65, 2.46, fill=LIGHT_GRAY, line="D4DBD6", rounded=True)
    text_box(slide, "10", 1.72, 3.03, 0.48, 0.45, size=24, color=DARK_GREEN, bold=True, align=PP_ALIGN.CENTER)
    text_box(slide, "localidades", 1.48, 3.48, 0.96, 0.23, size=10.2, color=DARK, align=PP_ALIGN.CENTER)
    # Dez pontos geográficos estilizados.
    points = [(0.92,3.22),(1.21,3.72),(0.87,4.15),(1.53,4.03),(2.44,3.26),(2.20,3.77),(2.62,4.25),(1.15,4.63),(1.82,4.56),(2.37,4.78)]
    for px, py in points:
        m = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(px), Inches(py), Inches(0.10), Inches(0.10))
        m.fill.solid(); m.fill.fore_color.rgb = rgb(GREEN); m.line.fill.background()
    text_box(slide, "dez séries temporais relacionadas", 1.10, 4.90, 1.70, 0.30, size=8.6, color=DARK, align=PP_ALIGN.CENTER)

    # Linha do tempo.
    timeline_y = 3.00
    add_line(slide, 4.20, timeline_y, 8.40, timeline_y, color=GREEN, width=2.0)
    timeline = [
        (4.30, "2019", "contexto\ninicial"),
        (6.10, "2020–2023", "48 alvos mensais de\ntreinamento por localidade"),
        (8.00, "2024", "12 previsões retrospectivas\npor localidade"),
    ]
    for x, year, body in timeline:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - 0.07), Inches(timeline_y - 0.07), Inches(0.14), Inches(0.14))
        dot.fill.solid(); dot.fill.fore_color.rgb = rgb(GREEN); dot.line.fill.background()
        text_box(slide, year, x - 0.48, 2.65, 0.96, 0.24, size=11.2, color=DARK_GREEN, bold=True, align=PP_ALIGN.CENTER)
        text_box(slide, body, x - 0.68, 3.20, 1.36, 0.52, size=8.2, color=DARK, align=PP_ALIGN.CENTER)

    text_box(slide, "Avaliação walk-forward", 3.55, 3.85, 2.20, 0.23, size=9.2, color=DARK_GREEN, bold=True)
    steps = ["histórico até\ndez/2023", "prever\njan/2024", "incorporar janeiro\nobservado", "prever\nfevereiro", "…", "prever\ndezembro"]
    step_x = [3.55, 4.48, 5.41, 6.43, 7.30, 7.72]
    widths = [0.72, 0.72, 0.82, 0.72, 0.24, 0.78]
    for x, w, label in zip(step_x, widths, steps, strict=True):
        if label == "…":
            text_box(slide, label, x, 4.35, w, 0.24, size=15, color=DARK_GREEN, bold=True, align=PP_ALIGN.CENTER)
        else:
            rect(slide, x, 4.20, w, 0.78, fill=WHITE, line="D2D9D5", rounded=True)
            text_box(slide, label, x + 0.04, 4.30, w - 0.08, 0.48, size=7.3, color=DARK, align=PP_ALIGN.CENTER)
    for x in (4.30, 5.23, 6.27, 7.18, 7.55): add_arrow(slide, x, 4.60, x + 0.13, 4.60, color=GREEN, width=1.2)

    # Faixa inferior mensal.
    monthly_items = ["contexto\n12 meses", "horizonte\n1 mês", "avaliação\nwalk-forward", "parâmetros fixos\nno teste", "sementes\n11, 23, 42, 67, 89"]
    for i, label in enumerate(monthly_items):
        x = 0.62 + i * 1.58
        rect(slide, x, 5.54, 1.42, 0.68, fill=LIGHT_GRAY, line="D6DDD8", rounded=True)
        text_box(slide, label, x + 0.08, 5.66, 1.26, 0.40, size=7.8, color=DARK, bold=i == 4, align=PP_ALIGN.CENTER)

    # Experimento horário.
    rect(slide, 9.12, 2.00, 3.55, 4.52, fill=WHITE, line=GREEN, rounded=True)
    text_box(slide, "EXPERIMENTO HORÁRIO", 9.48, 2.18, 2.78, 0.28, size=14, color=BLACK, bold=True)
    hourly = [
        "contexto: 336 horas",
        "saída direta: 72 horas",
        "horizontes: 24 h, 48 h e 72 h",
        "treino: 2019–2022",
        "validação: 2023",
        "teste: 2024",
    ]
    for i, label in enumerate(hourly):
        y = 2.78 + i * 0.55
        rect(slide, 9.38, y, 3.03, 0.42, fill=LIGHT_GRAY, line="D6DDD8", rounded=True)
        text_box(slide, label, 9.66, y + 0.04, 2.47, 0.34, size=9.3, color=DARK, bold=True)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.51), Inches(y + 0.14), Inches(0.12), Inches(0.12))
        dot.fill.solid(); dot.fill.fore_color.rgb = rgb(GREEN); dot.line.fill.background()
    rect(slide, 9.12, 6.02, 3.55, 0.40, fill=WHITE, line=GREEN, rounded=True)
    text_box(slide, "↔  Experimento horário separado do mensal", 9.32, 6.08, 3.15, 0.25, size=8.8, color=DARK_GREEN, bold=True, align=PP_ALIGN.CENTER)

    set_notes(slide, """
Explicar que existem dois experimentos separados.
No mensal, há dez séries relacionadas, contexto de 12 meses, horizonte de um mês, 48 alvos de treinamento por localidade entre 2020 e 2023 e 12 previsões retrospectivas em 2024.
O teste é walk-forward: o histórico observado cresce, mas os parâmetros do modelo permanecem fixos.
No horário, o contexto é de 336 horas e a saída direta contém 72 horas, avaliadas em 24, 48 e 72 horas. O treino usa 2019 a 2022, a validação usa 2023 e o teste usa 2024.
Transição: voltar à visão geral das famílias de modelos.
""")
    return slide


def add_intro_slides(prs: Presentation) -> None:
    slides = [add_need_slide(prs), add_ghi_slide(prs), add_protocol_slide(prs)]
    slide_ids = prs.slides._sldIdLst
    ids = list(slide_ids[-3:])
    for slide_id in ids:
        slide_ids.remove(slide_id)
    for position, slide_id in enumerate(ids, start=1):
        slide_ids.insert(position, slide_id)


def update_notes(prs: Presentation) -> None:
    notes = {
        1: """
Abrir apresentando o Share Knowledge: uma visão comparativa de modelos de previsão de séries temporais, indo de referências simples a redes recorrentes, modelos probabilísticos e arquiteturas baseadas em periodicidade.
Explicar que o GHI será usado apenas como caso prático para conectar os conceitos aos resultados.
Transição: começar pelo formato geral de uma tarefa de previsão.
""",
        2: """
Explicar o problema central: usar uma janela histórica para prever o próximo valor.
O gráfico é conceitual e não representa uma escala física; no experimento real, o alvo é a GHI média mensal em W/m².
Destacar o protocolo mensal: dez séries relacionadas, contexto de 12 meses, horizonte de um mês e teste retrospectivo em 2024.
Transição: mostrar que arquiteturas diferentes representam esse mesmo histórico de formas distintas.
""",
        3: """
Usar esta lâmina como mapa da apresentação.
Modelos tabulares usam atributos derivados; recorrentes processam sequências; probabilísticos retornam distribuições; TimesNet explora periodicidades em resolução horária.
As referências simples são indispensáveis para verificar se a complexidade realmente agrega valor.
O experimento do TimesNet já foi concluído, mas permanece separado do ranking mensal.
""",
        4: """
Explicar que XGBoost e MLP tratam a previsão como um problema tabular.
O XGBoost combina árvores de regressão sequencialmente; a MLP aprende relações não lineares entre os atributos.
Ambos dependem de uma boa engenharia de atributos temporais: 12 defasagens, médias de 3, 6 e 12 meses, calendário e identidade da série.
""",
        5: """
Apresentar a RNN como o primeiro modelo em que a ordem temporal faz parte da arquitetura.
Cada observação atualiza um estado oculto que funciona como resumo aprendido do passado; ele não contém necessariamente todo o histórico.
Informações distantes podem se perder durante o treinamento.
""",
        6: """
Explicar a LSTM como evolução da RNN para controlar melhor a memória.
As portas de esquecimento, entrada e saída regulam o fluxo de informação.
A arquitetura facilita dependências mais longas, mas não garante que toda informação relevante será preservada.
""",
        7: """
Apresentar a DilatedRNN como recorrência multiescala.
A execução utiliza dilatações 1, 2 e 4 para combinar dinâmica local e dependências mais espaçadas.
O campo temporal é ampliado, mas qualquer vantagem precisa ser confirmada pelos resultados.
""",
        8: """
Explicar que DeepAR muda a saída: em vez de apenas um número, prevê uma distribuição de futuros possíveis.
O modelo é global, aprende com as dez séries e utiliza duas camadas recorrentes de 40 unidades, contexto de 12 meses e amostras para formar quantis.
Uma faixa estreita não é automaticamente boa: também é preciso verificar cobertura e calibração.
""",
        9: """
Explicar DeepNPTS como abordagem probabilística não paramétrica.
A figura é didática; no experimento real, os candidatos correspondem às 12 posições do contexto mensal, de t−11 a t.
O modelo atribui probabilidades aos valores do contexto e produz uma distribuição discreta.
Ponto forte: suporte histórico plausível. Limitação: os resultados possíveis ficam restritos ao contexto disponível.
""",
        10: """
Apresentar TimesNet como experimento horário multietapas já concluído.
A arquitetura descobre períodos dominantes por FFT, reorganiza o sinal 1D em estruturas 2D e aplica convoluções para capturar padrões intra e interperíodo.
O protocolo usa 336 horas de contexto e avalia previsões diretas em 24, 48 e 72 horas.
Os resultados devem ser apresentados separadamente do ranking mensal.
""",
        11: """
Alinhar como o desempenho é medido antes de mostrar resultados.
O macro-MAE é o critério mensal principal e dá o mesmo peso às dez localidades.
RMSE penaliza erros grandes; nRMSE é o RMSE dividido pela média observada da localidade; R² é complementar.
Para previsões probabilísticas, ler CRPS, cobertura e largura conjuntamente.
""",
        12: """
Ler o gráfico pelo macro-MAE, em que menor é melhor.
A climatologia teve o menor valor observado, 12,07 W/m²; a LSTM foi o melhor modelo aprendido, com 12,22 W/m².
A diferença LSTM menos climatologia foi de apenas +0,15 W/m², com IC95% incluindo zero e resultado estatisticamente inconclusivo.
Os recorrentes ficaram próximos do topo, mas isso é uma observação descritiva, não prova de superioridade universal.
""",
    }
    for number, content in notes.items():
        set_notes(prs.slides[number - 1], content)

    # Depois da inserção, o slide final original passa à posição 14.
    final_slide = prs.slides[13]
    set_notes(
        final_slide,
        """
Comparar DeepAR e DeepNPTS pela combinação entre CRPS, cobertura e largura.
O DeepAR teve CRPS menor e intervalos mais estreitos, mas cobertura de apenas 61,67% para um intervalo nominal de 90%.
O DeepNPTS chegou a 90,83% de cobertura, porém com intervalos de 145,56 W/m², aproximadamente 4,1 vezes mais largos.
Fechar com as conclusões: baselines fortes são essenciais, complexidade não garante melhor desempenho e não existe um modelo universalmente vencedor.
Como próximos passos, repetir o experimento horário com múltiplas sementes, investigar variáveis meteorológicas adicionais e melhorar a calibração probabilística.
""",
    )


def fix_final_page_number(prs: Presentation) -> None:
    hourly_slide = prs.slides[-2]
    cover_and_text(
        hourly_slide,
        "15",
        12.61,
        0.12,
        0.53,
        0.34,
        size=10.5,
        color=DARK_GREEN,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )
    final_slide = prs.slides[-1]
    cover_and_text(
        final_slide,
        "16",
        12.61,
        0.12,
        0.53,
        0.34,
        size=10.5,
        color=DARK_GREEN,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )


def main() -> None:
    prs = Presentation(SOURCE)
    correct_existing_slides(prs)
    add_hourly_results_slide(prs)
    update_notes(prs)
    add_intro_slides(prs)
    fix_final_page_number(prs)
    prs.save(OUTPUT)
    print(OUTPUT)
    print(f"slides={len(prs.slides)}")


if __name__ == "__main__":
    main()
