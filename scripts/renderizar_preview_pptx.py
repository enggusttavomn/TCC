#!/usr/bin/env python3
"""Renderizador leve para inspeção visual de apresentações sem LibreOffice.

Não substitui o PowerPoint: cobre apenas os elementos usados pelo deck deste
projeto (retângulos, círculos, triângulos, linhas e caixas de texto).
"""

from __future__ import annotations

import math
import sys
import textwrap
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN


DPI = 120
EMU_PER_INCH = 914400
REGULAR_FONT = "/usr/share/fonts/truetype/lato/Lato-Medium.ttf"
BOLD_FONT = "/usr/share/fonts/truetype/lato/Lato-Heavy.ttf"


def px(emu: int) -> int:
    return round(emu / EMU_PER_INCH * DPI)


def color_from(obj, default=None):
    try:
        value = obj.rgb
        if value is not None:
            return tuple(value)
    except (AttributeError, TypeError, ValueError):
        pass
    return default


def shape_fill(shape):
    try:
        if shape.fill.type is not None:
            return color_from(shape.fill.fore_color, None)
    except (AttributeError, TypeError, ValueError):
        pass
    return None


def shape_line(shape):
    try:
        return color_from(shape.line.color, None)
    except (AttributeError, TypeError, ValueError):
        return None


def get_font(size_pt: float, bold: bool):
    font_path = BOLD_FONT if bold else REGULAR_FONT
    return ImageFont.truetype(font_path, max(7, round(size_pt * DPI / 72)))


def wrap_text(draw, text: str, font, max_width: int):
    lines: list[str] = []
    for raw in text.splitlines() or [""]:
        if not raw:
            lines.append("")
            continue
        words = raw.split()
        line = ""
        for word in words:
            candidate = f"{line} {word}".strip()
            if draw.textlength(candidate, font=font) <= max_width or not line:
                line = candidate
            else:
                lines.append(line)
                line = word
        lines.append(line)
    return lines


def draw_text_frame(canvas, draw, shape, bbox):
    if not shape.has_text_frame or not shape.text.strip():
        return
    tf = shape.text_frame
    left, top, right, bottom = bbox
    left += px(tf.margin_left)
    right -= px(tf.margin_right)
    top += px(tf.margin_top)
    bottom -= px(tf.margin_bottom)
    available_w = max(1, right - left)
    available_h = max(1, bottom - top)

    blocks = []
    for paragraph in tf.paragraphs:
        if paragraph.runs:
            run = paragraph.runs[0]
            size = run.font.size.pt if run.font.size else 18
            bold = bool(run.font.bold)
            color = color_from(run.font.color, (245, 248, 252))
            text = "".join(r.text for r in paragraph.runs)
        else:
            size, bold, color, text = 18, False, (245, 248, 252), paragraph.text
        font = get_font(size, bold)
        lines = wrap_text(draw, text, font, available_w)
        line_h = max(1, round(font.size * 1.14))
        blocks.append((paragraph, font, color, lines, line_h))

    total_h = sum(len(lines) * line_h for _, _, _, lines, line_h in blocks)
    if tf.vertical_anchor == MSO_ANCHOR.MIDDLE:
        y = top + max(0, (available_h - total_h) // 2)
    elif tf.vertical_anchor == MSO_ANCHOR.BOTTOM:
        y = bottom - total_h
    else:
        y = top

    for paragraph, font, color, lines, line_h in blocks:
        for line in lines:
            line_w = draw.textlength(line, font=font)
            if paragraph.alignment == PP_ALIGN.CENTER:
                x = left + max(0, (available_w - line_w) / 2)
            elif paragraph.alignment == PP_ALIGN.RIGHT:
                x = right - line_w
            else:
                x = left
            draw.text((x, y), line, font=font, fill=color)
            y += line_h


def draw_auto_shape(draw, shape, bbox):
    fill = shape_fill(shape)
    outline = shape_line(shape)
    width = max(1, round((shape.line.width.pt if shape.line.width else 1) * DPI / 72))
    left, top, right, bottom = bbox
    kind = shape.auto_shape_type

    if kind == MSO_AUTO_SHAPE_TYPE.OVAL:
        draw.ellipse(bbox, fill=fill, outline=outline, width=width)
    elif kind == MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE:
        radius = max(2, min(right - left, bottom - top) // 8)
        draw.rounded_rectangle(bbox, radius=radius, fill=fill, outline=outline, width=width)
    elif kind == MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE:
        points = [
            ((left + right) / 2, top),
            (right, bottom),
            (left, bottom),
        ]
        angle = math.radians(shape.rotation or 0)
        if angle:
            cx, cy = (left + right) / 2, (top + bottom) / 2
            rotated = []
            for x, y in points:
                dx, dy = x - cx, y - cy
                rotated.append(
                    (
                        cx + dx * math.cos(angle) - dy * math.sin(angle),
                        cy + dx * math.sin(angle) + dy * math.cos(angle),
                    )
                )
            points = rotated
        draw.polygon(points, fill=fill, outline=outline)
    else:
        draw.rectangle(bbox, fill=fill, outline=outline, width=width)


def render_slide(prs, slide):
    canvas = Image.new(
        "RGB",
        (px(prs.slide_width), px(prs.slide_height)),
        color_from(slide.background.fill.fore_color, (7, 20, 33)),
    )
    draw = ImageDraw.Draw(canvas)
    for shape in slide.shapes:
        left, top = px(shape.left), px(shape.top)
        right, bottom = left + px(shape.width), top + px(shape.height)
        bbox = (left, top, right, bottom)
        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            color = shape_line(shape) or (39, 70, 92)
            width = max(
                1,
                round((shape.line.width.pt if shape.line.width else 1) * DPI / 72),
            )
            draw.line((left, top, right, bottom), fill=color, width=width)
        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            draw_auto_shape(draw, shape, bbox)
        draw_text_frame(canvas, draw, shape, bbox)
    return canvas


def main():
    source = (
        Path(sys.argv[1])
        if len(sys.argv) > 1
        else Path("apresentacao_modelos_machine_learning.pptx")
    )
    destination = source.with_name(source.stem + "_preview.png")
    prs = Presentation(source)
    slides = [render_slide(prs, slide) for slide in prs.slides]

    thumb_w, thumb_h = 480, 270
    gap = 18
    cols = 3
    rows = math.ceil(len(slides) / cols)
    sheet = Image.new(
        "RGB",
        (
            cols * thumb_w + (cols + 1) * gap,
            rows * thumb_h + (rows + 1) * gap,
        ),
        (3, 10, 17),
    )
    for i, slide in enumerate(slides):
        thumb = slide.resize((thumb_w, thumb_h), Image.Resampling.LANCZOS)
        x = gap + (i % cols) * (thumb_w + gap)
        y = gap + (i // cols) * (thumb_h + gap)
        sheet.paste(thumb, (x, y))
    sheet.save(destination, quality=95)
    print(f"Prévia criada: {destination}")
    print(f"Dimensões: {sheet.width}x{sheet.height}")


if __name__ == "__main__":
    main()
