#!/usr/bin/env python3
"""Generate the Machine Learning models deck using the SiDi template."""

from __future__ import annotations

import math
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / ".backup_apresentacao" / "2026_SiDi Automation_Ai.pptx"
OUTPUT = ROOT / "00_APRESENTACAO" / "apresentacao_modelos_machine_learning.pptx"

W, H = 10.0, 5.625
FONT = "Arial"

BLACK = "111111"
DARK = "272727"
GRAY = "676D6A"
MID = "A2A8A5"
LIGHT = "F2F4F2"
LIGHT_2 = "E7ECE9"
WHITE = "FFFFFF"
GREEN = "00B83F"
DARK_GREEN = "006325"
MINT = "CFF7DA"
PURPLE = "6F42C1"
BLUE = "0061D1"
PINK = "F55496"
YELLOW = "F7C743"
RED = "FF555D"
ORANGE = "F28C28"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def clear_existing_slides(prs: Presentation) -> None:
    ids = prs.slides._sldIdLst
    for slide_id in list(ids):
        prs.part.drop_rel(slide_id.rId)
        ids.remove(slide_id)


def add_shape(
    slide,
    kind,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str | None = None,
    line: str | None = None,
    line_width: float = 1.0,
):
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 16,
    color: str = BLACK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    p.line_spacing = 1.0
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_line(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    *,
    color: str = MID,
    width: float = 1.2,
):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    return line


def add_arrowhead(slide, x: float, y: float, color: str, size: float = 0.10):
    tri = add_shape(
        slide,
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        x,
        y,
        size,
        size,
        fill=color,
    )
    tri.rotation = 90


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = LIGHT,
    line: str | None = None,
    accent: str | None = None,
):
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x,
        y,
        w,
        h,
        fill=fill,
        line=line,
    )
    if accent:
        add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x,
            y,
            0.07,
            h,
            fill=accent,
        )


def add_pill(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    *,
    fill: str = MINT,
    color: str = DARK_GREEN,
    line: str | None = None,
    size: float = 9.5,
):
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x,
        y,
        w,
        0.30,
        fill=fill,
        line=line,
    )
    add_text(
        slide,
        text,
        x,
        y + 0.01,
        w,
        0.25,
        size=size,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_bullet(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    *,
    dot: str = GREEN,
    size: float = 11.5,
    color: str = DARK,
):
    add_shape(slide, MSO_SHAPE.OVAL, x, y + 0.10, 0.07, 0.07, fill=dot)
    add_text(slide, text, x + 0.17, y, w - 0.17, 0.34, size=size, color=color)


def content_slide(prs: Presentation, kicker: str, title: str, number: int):
    slide = prs.slides.add_slide(prs.slide_layouts[11])
    add_text(slide, kicker.upper(), 0.52, 0.23, 4.4, 0.18, size=8.5, color=DARK_GREEN, bold=True)
    add_text(slide, title, 0.52, 0.48, 8.62, 0.46, size=23, color=BLACK, bold=True)
    add_text(
        slide,
        f"{number:02d}",
        9.20,
        0.29,
        0.28,
        0.18,
        size=8.5,
        color=MID,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )
    add_line(slide, 0.52, 1.05, 9.48, 1.05, color=GREEN, width=1.4)
    return slide


def set_placeholder(
    placeholder,
    text: str,
    *,
    size: float,
    color: str = BLACK,
    bold: bool = False,
):
    tf = placeholder.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = rgb(color)
    run.font.bold = bold


def slide_01(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.placeholders[10]
    subtitle = slide.placeholders[11]
    set_placeholder(title, "Machine Learning Models", size=28, bold=False)
    set_placeholder(
        subtitle,
        "Share Knowledge\nXGBoost · MLP · RNN · LSTM · DilatedRNN · DeepAR · DeepNPTS-inspired",
        size=12.5,
        color=GRAY,
    )
    return slide


def slide_02(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_placeholder(
        slide.placeholders[0],
        "Seven models,\nthree design choices",
        size=27,
        bold=True,
    )
    set_placeholder(
        slide.placeholders[1],
        "Representation  ·  Memory  ·  Predictive output",
        size=13,
        color=GRAY,
    )
    return slide


def draw_trees(slide, x: float, y: float):
    colors = [ORANGE, YELLOW, PINK]
    for i, color in enumerate(colors):
        ox = x + i * 0.78
        add_line(slide, ox + 0.28, y + 0.12, ox + 0.12, y + 0.43, color=GRAY)
        add_line(slide, ox + 0.28, y + 0.12, ox + 0.44, y + 0.43, color=GRAY)
        add_shape(slide, MSO_SHAPE.OVAL, ox + 0.20, y + 0.04, 0.17, 0.17, fill=color)
        add_shape(slide, MSO_SHAPE.OVAL, ox + 0.05, y + 0.38, 0.15, 0.15, fill=WHITE, line=color)
        add_shape(slide, MSO_SHAPE.OVAL, ox + 0.37, y + 0.38, 0.15, 0.15, fill=WHITE, line=color)
        add_shape(slide, MSO_SHAPE.OVAL, ox + 0.05, y + 0.68, 0.12, 0.12, fill=color)
        add_shape(slide, MSO_SHAPE.OVAL, ox + 0.39, y + 0.68, 0.12, 0.12, fill=color)


def draw_mlp(slide, x: float, y: float):
    layers = [(x, 4, BLUE), (x + 0.75, 5, GREEN), (x + 1.50, 2, PURPLE)]
    coords = []
    for lx, count, color in layers:
        ys = [y + i * (0.92 / (count - 1)) for i in range(count)]
        layer = []
        for yy in ys:
            layer.append((lx, yy))
            add_shape(slide, MSO_SHAPE.OVAL, lx, yy, 0.14, 0.14, fill=color)
        coords.append(layer)
    for left, right in zip(coords[:-1], coords[1:]):
        for x1, y1 in left:
            for x2, y2 in right:
                add_line(slide, x1 + 0.14, y1 + 0.07, x2, y2 + 0.07, color=LIGHT_2, width=0.5)


def slide_03(prs):
    slide = content_slide(prs, "Tabular learners", "When time has already been engineered into features", 3)
    add_card(slide, 0.52, 1.32, 4.28, 3.75, fill=LIGHT, accent=ORANGE)
    add_pill(slide, "XGBoost", 0.82, 1.58, 1.18, fill=WHITE, color=ORANGE, line=ORANGE)
    add_text(
        slide,
        "Each boosting round corrects\nresidual errors from the last.",
        0.82,
        2.04,
        2.80,
        0.58,
        size=16,
        bold=True,
    )
    draw_trees(slide, 1.02, 2.82)
    add_text(slide, "BOOSTING ROUNDS", 1.02, 3.67, 2.18, 0.14, size=7.5, color=GRAY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Σ", 3.58, 3.03, 0.36, 0.30, size=22, color=ORANGE, bold=True)
    add_shape(slide, MSO_SHAPE.OVAL, 4.10, 3.08, 0.22, 0.22, fill=ORANGE)
    add_bullet(slide, "Strong baseline for tabular data.", 0.82, 4.02, 3.35, dot=ORANGE)
    add_bullet(slide, "Needs lags, rolling statistics and calendar.", 0.82, 4.42, 3.35, dot=ORANGE, size=10.5)

    add_card(slide, 5.20, 1.32, 4.28, 3.75, fill=LIGHT, accent=BLUE)
    add_pill(slide, "MLP", 5.50, 1.58, 0.88, fill=WHITE, color=BLUE, line=BLUE)
    add_text(
        slide,
        "Dense layers learn smooth,\nnonlinear interactions.",
        5.50,
        2.04,
        2.65,
        0.58,
        size=16,
        bold=True,
    )
    draw_mlp(slide, 5.88, 2.82)
    add_text(slide, "INPUT", 5.74, 3.86, 0.42, 0.14, size=7, color=GRAY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "HIDDEN", 6.44, 3.86, 0.56, 0.14, size=7, color=GRAY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "OUTPUT", 7.22, 3.86, 0.55, 0.14, size=7, color=GRAY, bold=True, align=PP_ALIGN.CENTER)
    add_line(slide, 7.56, 3.36, 8.60, 3.36, color=PURPLE, width=1.8)
    add_arrowhead(slide, 8.55, 3.31, PURPLE)
    add_shape(slide, MSO_SHAPE.OVAL, 8.77, 3.25, 0.22, 0.22, fill=PURPLE)
    add_bullet(slide, "Flexible nonlinear function approximator.", 5.50, 4.02, 3.55, dot=BLUE, size=10.5)
    add_bullet(slide, "Requires scaling and regularization.", 5.50, 4.42, 3.55, dot=BLUE, size=10.5)
    return slide


def slide_04(prs):
    slide = content_slide(prs, "Sequence models", "RNN: the hidden state carries context", 4)
    add_text(
        slide,
        "Each step combines the current observation with a compressed summary of the past.",
        0.52,
        1.21,
        6.5,
        0.28,
        size=12,
        color=GRAY,
    )
    add_card(slide, 7.42, 1.16, 2.06, 0.43, fill=MINT)
    add_text(slide, "hₜ = φ(Wₓxₜ + Wₕhₜ₋₁)", 7.57, 1.30, 1.76, 0.17, size=10.5, color=DARK_GREEN, bold=True, align=PP_ALIGN.CENTER)
    y = 2.45
    for i in range(5):
        x = 0.63 + i * 1.75
        label = f"xₜ₋{4-i}" if i < 4 else "xₜ"
        add_pill(slide, label, x + 0.26, 1.70, 0.62, fill=MINT, color=DARK_GREEN, size=10)
        add_line(slide, x + 0.57, 2.00, x + 0.57, y, color=GREEN, width=1.3)
        add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x,
            y,
            1.15,
            0.78,
            fill=LIGHT,
            line=GREEN if i == 4 else LIGHT_2,
        )
        add_text(slide, "RNN", x, y + 0.16, 1.15, 0.23, size=13, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, f"h{i+1}", x, y + 0.48, 1.15, 0.16, size=8, color=GRAY, align=PP_ALIGN.CENTER)
        if i < 4:
            add_line(slide, x + 1.15, y + 0.39, x + 1.66, y + 0.39, color=ORANGE, width=1.5)
            add_arrowhead(slide, x + 1.61, y + 0.34, ORANGE)
    add_line(slide, 7.63, 3.23, 7.63, 3.75, color=GREEN, width=1.5)
    add_arrowhead(slide, 7.58, 3.70, GREEN)
    add_pill(slide, "forecast ŷₜ₊₁", 7.02, 3.85, 1.30, fill=MINT, color=DARK_GREEN, size=10)
    add_card(slide, 0.63, 4.28, 8.70, 0.72, fill=WHITE, line=LIGHT_2, accent=GREEN)
    add_text(slide, "STRENGTH", 0.91, 4.47, 0.70, 0.18, size=8.5, color=DARK_GREEN, bold=True)
    add_text(slide, "native ordering and local dependencies", 1.72, 4.43, 2.58, 0.24, size=10.5, bold=True)
    add_text(slide, "LIMIT", 4.66, 4.47, 0.48, 0.18, size=8.5, color=RED, bold=True)
    add_text(slide, "distant signals may fade during training", 5.27, 4.43, 3.55, 0.24, size=10.5, color=GRAY)
    return slide


def slide_05(prs):
    slide = content_slide(prs, "Selective memory", "LSTM: remembering also means learning what to forget", 5)
    add_card(slide, 0.52, 1.35, 6.55, 3.53, fill=LIGHT)
    add_text(slide, "CELL-STATE HIGHWAY  cₜ₋₁", 0.86, 1.64, 2.10, 0.20, size=8.5, color=GRAY, bold=True)
    add_line(slide, 0.88, 2.11, 6.69, 2.11, color=GREEN, width=3.0)
    add_arrowhead(slide, 6.60, 2.04, GREEN, size=0.14)
    add_text(slide, "cₜ", 6.34, 1.63, 0.32, 0.20, size=10, color=DARK_GREEN, bold=True)
    gates = [
        (1.20, "FORGET GATE", "discard noise", ORANGE, "×"),
        (3.00, "INPUT GATE", "write evidence", GREEN, "+"),
        (4.80, "OUTPUT GATE", "expose hₜ", PURPLE, "σ"),
    ]
    for x, label, sub, color, symbol in gates:
        add_shape(slide, MSO_SHAPE.OVAL, x, 2.66, 0.52, 0.52, fill=WHITE, line=color, line_width=1.8)
        add_text(slide, symbol, x, 2.77, 0.52, 0.22, size=15, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_line(slide, x + 0.26, 2.66, x + 0.26, 2.14, color=color, width=1.4)
        add_text(slide, label, x - 0.25, 3.42, 1.02, 0.18, size=8, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, sub, x - 0.36, 3.72, 1.24, 0.22, size=9.5, color=GRAY, align=PP_ALIGN.CENTER)
    add_pill(slide, "xₜ + hₜ₋₁", 2.78, 4.27, 1.34, fill=WHITE, color=BLACK, line=LIGHT_2, size=10)

    add_card(slide, 7.35, 1.35, 2.13, 3.53, fill=WHITE, line=LIGHT_2, accent=GREEN)
    add_text(slide, "WHY IT WORKS", 7.68, 1.72, 1.45, 0.20, size=8.5, color=DARK_GREEN, bold=True)
    add_text(
        slide,
        "Controlled memory keeps useful signals alive across longer horizons.",
        7.68,
        2.15,
        1.44,
        1.10,
        size=13.5,
        bold=True,
    )
    add_text(slide, "TRADE-OFF", 7.68, 3.57, 0.84, 0.20, size=8.5, color=ORANGE, bold=True)
    add_text(
        slide,
        "More parameters and a more demanding optimization process.",
        7.68,
        3.90,
        1.42,
        0.62,
        size=10.5,
        color=GRAY,
    )
    return slide


def slide_06(prs):
    slide = content_slide(prs, "Multi-scale recurrence", "DilatedRNN: structured jumps across time", 6)
    add_text(
        slide,
        "A dilation s connects the current state directly to hₜ₋ₛ, creating a multi-resolution temporal hierarchy.",
        0.52,
        1.21,
        6.7,
        0.24,
        size=11.2,
        color=GRAY,
    )
    add_card(slide, 7.22, 1.18, 2.26, 0.48, fill=MINT)
    add_text(
        slide,
        "hₜ = φ(Wₓxₜ + Wₕhₜ₋ₛ + b)",
        7.38,
        1.34,
        1.94,
        0.18,
        size=10.5,
        color=DARK_GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    rows = [
        (2.06, 1, GREEN, "LOCAL", "adjacent dynamics"),
        (3.03, 2, ORANGE, "MID-RANGE", "periodic structure"),
        (4.00, 4, PURPLE, "LONG-RANGE", "coarse temporal context"),
    ]
    node_xs = [2.25 + i * 0.67 for i in range(11)]
    for yy, step, color, label, descriptor in rows:
        add_pill(slide, f"s = {step}", 0.70, yy - 0.05, 0.72, fill=WHITE, color=color, line=color, size=8.5)
        add_text(slide, label, 0.70, yy + 0.40, 1.10, 0.15, size=7.5, color=color, bold=True)
        add_text(slide, descriptor, 0.70, yy + 0.58, 1.35, 0.16, size=8, color=GRAY)
        for target in range(step, len(node_xs)):
            x1 = node_xs[target - step] + 0.12
            x2 = node_xs[target] + 0.12
            add_line(slide, x1, yy + 0.12, x2, yy + 0.12, color=color, width=1.25)
        for i, x in enumerate(node_xs):
            active = i >= step
            add_shape(
                slide,
                MSO_SHAPE.OVAL,
                x,
                yy,
                0.24,
                0.24,
                fill=color if i == len(node_xs) - 1 else WHITE,
                line=color if active else LIGHT_2,
                line_width=1.1,
            )
            if yy == rows[-1][0]:
                add_text(slide, str(i), x - 0.02, yy + 0.34, 0.28, 0.13, size=7, color=GRAY, align=PP_ALIGN.CENTER)
        add_text(slide, "hₜ", 9.30, yy + 0.02, 0.18, 0.14, size=8.5, color=color, bold=True)

    add_card(slide, 0.70, 5.00, 8.78, 0.32, fill=LIGHT)
    add_text(
        slide,
        "Implementation note  ·  true sequences enter the recurrent encoder; auxiliary covariates are fused afterwards.",
        0.88,
        5.10,
        8.42,
        0.14,
        size=8.3,
        color=GRAY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    return slide


def slide_07(prs):
    slide = content_slide(prs, "Probabilistic forecasting", "DeepAR: a distribution over possible futures", 7)
    add_card(slide, 0.52, 1.35, 2.52, 3.53, fill=LIGHT)
    add_text(slide, "RELATED TIME SERIES", 0.82, 1.65, 1.50, 0.18, size=8.5, color=DARK_GREEN, bold=True)
    colors = [GREEN, ORANGE, PURPLE, PINK]
    for row, color in enumerate(colors):
        points = []
        for i in range(8):
            xx = 0.82 + i * 0.24
            yy = 2.20 + row * 0.48 - 0.08 * math.sin(i * 1.2 + row)
            points.append((xx, yy))
        for p1, p2 in zip(points[:-1], points[1:]):
            add_line(slide, *p1, *p2, color=color, width=1.4)
    add_pill(slide, "shared global parameters", 0.82, 4.24, 1.82, fill=WHITE, color=DARK, line=LIGHT_2, size=8.3)
    add_line(slide, 3.04, 3.08, 3.66, 3.08, color=GREEN, width=1.7)
    add_arrowhead(slide, 3.61, 3.03, GREEN)
    add_shape(slide, MSO_SHAPE.OVAL, 3.84, 2.53, 1.10, 1.10, fill=WHITE, line=GREEN, line_width=1.8)
    add_text(slide, "RNN", 3.84, 2.76, 1.10, 0.24, size=15, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "global", 3.84, 3.08, 1.10, 0.16, size=8, color=GRAY, align=PP_ALIGN.CENTER)
    add_line(slide, 4.94, 3.08, 5.55, 3.08, color=GREEN, width=1.7)
    add_arrowhead(slide, 5.50, 3.03, GREEN)

    add_card(slide, 5.70, 1.35, 3.78, 3.53, fill=WHITE, line=LIGHT_2)
    add_pill(slide, "likelihood", 5.02, 2.66, 0.62, fill=WHITE, color=GRAY, line=LIGHT_2, size=7.5)
    add_text(slide, "PREDICTIVE DISTRIBUTION", 6.00, 1.65, 2.35, 0.18, size=8.5, color=PINK, bold=True)
    x0, y0 = 6.02, 4.33
    add_line(slide, x0, y0, 9.12, y0, color=LIGHT_2)
    add_line(slide, x0, 2.08, x0, y0, color=LIGHT_2)
    hist = []
    for i in range(7):
        xx = x0 + i * 0.27
        yy = 3.84 - 0.12 * i - 0.14 * math.sin(i * 0.8)
        hist.append((xx, yy))
    for p1, p2 in zip(hist[:-1], hist[1:]):
        add_line(slide, *p1, *p2, color=BLACK, width=1.5)
    split = hist[-1][0]
    add_line(slide, split, 2.05, split, y0, color=MID)
    med0 = hist[-1][1]
    upper, med, lower = [], [], []
    for i in range(1, 6):
        xx = split + i * 0.29
        center = med0 - 0.08 * i
        spread = 0.07 + 0.06 * i
        upper.append((xx, center - spread))
        med.append((xx, center))
        lower.append((xx, center + spread))
        add_line(slide, xx, center - spread, xx, center + spread, color="F9D7E7", width=6)
    for arr, color, width in [(upper, PINK, 1), (lower, PINK, 1), (med, DARK_GREEN, 2)]:
        for p1, p2 in zip([hist[-1]] + arr[:-1], arr):
            add_line(slide, *p1, *p2, color=color, width=width)
    add_pill(slide, "P10", 8.20, 2.20, 0.50, fill=LIGHT, color=PINK, size=8)
    add_pill(slide, "P50", 8.76, 2.20, 0.50, fill=MINT, color=DARK_GREEN, size=8)
    add_pill(slide, "P90", 8.20, 2.55, 0.50, fill=LIGHT, color=PINK, size=8)
    add_text(
        slide,
        "Decision value: calibrated uncertainty for inventory, capacity, risk and energy.",
        0.72,
        5.06,
        8.80,
        0.22,
        size=11,
        color=DARK_GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    return slide


def slide_08(prs):
    slide = content_slide(prs, "Experimental approach", "DeepNPTS-inspired: history as a candidate set", 8)
    add_text(
        slide,
        "The forecast reuses observed values, assigning probability mass through similarity and recency.",
        0.52,
        1.21,
        8.3,
        0.26,
        size=11.5,
        color=GRAY,
    )
    stages = [
        ("1", "CONTEXT", "recent pattern", GREEN),
        ("2", "SCORE", "similarity + recency", ORANGE),
        ("3", "WEIGHTS", "softmax over history", PURPLE),
        ("4", "FORECAST", "empirical distribution", YELLOW),
    ]
    for i, (num, label, sub, color) in enumerate(stages):
        x = 0.52 + i * 2.37
        add_card(slide, x, 1.70, 2.03, 2.84, fill=LIGHT)
        add_shape(slide, MSO_SHAPE.OVAL, x + 0.18, 1.91, 0.28, 0.28, fill=color)
        add_text(slide, num, x + 0.18, 1.97, 0.28, 0.14, size=8.5, color=WHITE if color != YELLOW else BLACK, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, label, x + 0.58, 1.97, 1.16, 0.16, size=8.5, color=color, bold=True)
        add_text(slide, sub, x + 0.18, 4.10, 1.67, 0.24, size=9, color=GRAY, align=PP_ALIGN.CENTER)
        if i < 3:
            add_line(slide, x + 2.03, 3.05, x + 2.28, 3.05, color=color, width=1.5)
            add_arrowhead(slide, x + 2.23, 3.00, color)
    points = []
    for i in range(7):
        points.append((0.83 + i * 0.20, 3.30 - 0.18 * math.sin(i)))
    for p1, p2 in zip(points[:-1], points[1:]):
        add_line(slide, *p1, *p2, color=GREEN, width=1.7)
    for i, width in enumerate([0.35, 0.72, 1.10, 0.55]):
        add_shape(slide, MSO_SHAPE.RECTANGLE, 3.05, 2.52 + i * 0.31, width, 0.13, fill=ORANGE)
    add_text(slide, "Σ pᵢ = 1", 5.45, 2.49, 1.10, 0.26, size=15, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
    for i, val in enumerate([0.14, 0.38, 0.22, 0.68, 0.90]):
        size = 0.10 + 0.18 * val
        add_shape(slide, MSO_SHAPE.OVAL, 5.28 + i * 0.29, 3.20 - size / 2, size, size, fill=PURPLE)
    for x, y, size in [(7.85, 3.50, 0.17), (8.22, 2.77, 0.26), (8.68, 3.18, 0.21), (9.00, 2.50, 0.32)]:
        add_line(slide, x + size / 2, 3.78, x + size / 2, y + size / 2, color=MID)
        add_shape(slide, MSO_SHAPE.OVAL, x, y, size, size, fill=YELLOW)
    add_line(slide, 7.78, 3.78, 9.36, 3.78, color=MID)
    add_card(slide, 0.52, 4.82, 8.96, 0.43, fill=MINT, accent=GREEN)
    add_text(
        slide,
        "Experimental approximation inspired by the non-parametric principle — not a full reproduction of the official DeepNPTS.",
        0.78,
        4.95,
        8.42,
        0.17,
        size=8.5,
        color=DARK_GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    return slide


def slide_09(prs):
    slide = content_slide(prs, "Comparison", "Seven models at a glance", 9)
    x0 = 0.52
    widths = [1.15, 1.42, 1.48, 1.12, 2.00, 1.79]
    headers = ["MODEL", "INPUT", "MEMORY", "OUTPUT", "STRENGTH", "WATCH FOR"]
    cur = x0
    for label, width in zip(headers, widths):
        add_shape(slide, MSO_SHAPE.RECTANGLE, cur, 1.32, width, 0.37, fill=DARK_GREEN)
        add_text(slide, label, cur + 0.09, 1.43, width - 0.18, 0.14, size=7.3, color=WHITE, bold=True)
        cur += width
    rows = [
        ("XGBoost", "features", "explicit lags", "point", "tabular data", "feature design", ORANGE),
        ("MLP", "features", "explicit lags", "point", "nonlinearity", "scaling", BLUE),
        ("RNN", "sequence", "short", "point", "temporal order", "long distance", GREEN),
        ("LSTM", "sequence", "selective", "point", "long memory", "parameter count", DARK_GREEN),
        ("DilatedRNN", "sequence", "multi-scale", "point", "mixed rhythms", "dilation design", PURPLE),
        ("DeepAR", "many series", "autoregr.", "distribution", "uncertainty", "likelihood choice", PINK),
        ("DeepNPTS*", "history", "similarity", "distribution", "plausibility", "historical support", YELLOW),
    ]
    row_h = 0.45
    for row_idx, row in enumerate(rows):
        y = 1.69 + row_idx * row_h
        fill = WHITE if row_idx % 2 else LIGHT
        cur = x0
        for col_idx, (value, width) in enumerate(zip(row[:-1], widths)):
            add_shape(slide, MSO_SHAPE.RECTANGLE, cur, y, width, row_h, fill=fill)
            add_text(
                slide,
                value,
                cur + 0.09,
                y + 0.14,
                width - 0.18,
                0.16,
                size=8.2,
                color=row[-1] if col_idx == 0 else (DARK if col_idx < 4 else GRAY),
                bold=col_idx == 0,
            )
            cur += width
        add_shape(slide, MSO_SHAPE.RECTANGLE, x0, y, 0.04, row_h, fill=row[-1])
    add_text(slide, "* experimental approximation inspired by DeepNPTS", 0.54, 5.02, 3.65, 0.16, size=8, color=GRAY)
    add_pill(slide, "point = best estimate", 5.42, 4.94, 1.66, fill=LIGHT, color=DARK, size=8.3)
    add_pill(slide, "distribution = scenarios + risk", 7.22, 4.94, 2.22, fill=MINT, color=DARK_GREEN, size=8.3)
    return slide


def slide_10(prs):
    slide = content_slide(prs, "Model selection", "Choose by capability, not by popularity", 10)
    add_text(
        slide,
        "Map the decision requirement first. Benchmark candidates inside the relevant region.",
        0.52,
        1.21,
        8.8,
        0.24,
        size=11.5,
        color=GRAY,
    )
    add_card(slide, 0.72, 1.58, 7.36, 3.42, fill=LIGHT)
    # Capability matrix.
    add_line(slide, 1.42, 4.48, 7.72, 4.48, color=DARK, width=1.1)
    add_arrowhead(slide, 7.66, 4.43, DARK)
    add_line(slide, 1.42, 4.48, 1.42, 1.92, color=DARK, width=1.1)
    up = add_shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, 1.37, 1.82, 0.10, 0.10, fill=DARK)
    up.rotation = 0
    add_line(slide, 4.57, 1.90, 4.57, 4.48, color=WHITE, width=1.0)
    add_line(slide, 1.42, 3.18, 7.72, 3.18, color=WHITE, width=1.0)
    add_text(slide, "ENGINEERED FEATURES", 1.55, 4.67, 2.05, 0.16, size=7.5, color=GRAY, bold=True)
    add_text(slide, "NATIVE TEMPORAL STRUCTURE", 5.50, 4.67, 2.18, 0.16, size=7.5, color=GRAY, bold=True, align=PP_ALIGN.RIGHT)
    y_label = add_text(slide, "UNCERTAINTY REQUIRED", 0.84, 2.27, 1.62, 0.17, size=7.5, color=GRAY, bold=True)
    y_label.rotation = 270
    add_text(slide, "LOW", 1.03, 4.25, 0.28, 0.14, size=7, color=GRAY, bold=True)
    add_text(slide, "HIGH", 0.99, 1.95, 0.34, 0.14, size=7, color=GRAY, bold=True)

    placements = [
        ("XGBoost", 1.82, 3.87, 1.12, ORANGE),
        ("MLP", 3.10, 3.58, 0.72, BLUE),
        ("RNN", 5.10, 3.91, 0.72, GREEN),
        ("LSTM", 5.98, 3.54, 0.82, DARK_GREEN),
        ("DilatedRNN", 6.42, 2.88, 1.18, PURPLE),
        ("DeepNPTS*", 3.53, 2.35, 1.24, YELLOW),
        ("DeepAR", 6.13, 2.04, 0.98, PINK),
    ]
    for label, x, y, width, color in placements:
        add_pill(slide, label, x, y, width, fill=WHITE, color=color, line=color, size=8.2)

    add_card(slide, 8.32, 1.58, 1.16, 3.42, fill=WHITE, line=LIGHT_2, accent=GREEN)
    add_text(slide, "START HERE", 8.56, 1.90, 0.68, 0.18, size=8.5, color=DARK_GREEN, bold=True, align=PP_ALIGN.CENTER)
    for i, (num, text) in enumerate(
        [
            ("1", "Define the output"),
            ("2", "Build a baseline"),
            ("3", "Earn added complexity"),
        ]
    ):
        yy = 2.36 + i * 0.78
        add_shape(slide, MSO_SHAPE.OVAL, 8.44, yy, 0.22, 0.22, fill=MINT)
        add_text(slide, num, 8.44, yy + 0.05, 0.22, 0.12, size=7.5, color=DARK_GREEN, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, text, 8.72, yy - 0.01, 0.62, 0.43, size=7.8, color=DARK, bold=True)
    add_text(
        slide,
        "The map narrows candidates; time-aware validation still decides.",
        0.72,
        5.14,
        8.76,
        0.18,
        size=9.5,
        color=DARK_GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    return slide


def slide_11(prs):
    slide = content_slide(prs, "Evaluation", "A fair experiment matters more than architecture", 11)
    add_text(slide, "Freeze the protocol; vary only the model.", 0.52, 1.21, 6.1, 0.24, size=11.5, color=GRAY)
    steps = [
        ("01", "TIME-BASED SPLIT", "train before test", GREEN),
        ("02", "SAME INPUTS", "comparable features", ORANGE),
        ("03", "SAME BUDGET", "controlled search", PURPLE),
        ("04", "RIGHT METRICS", "error + calibration", DARK_GREEN),
        ("05", "REPLICATION", "seeds + intervals", PINK),
    ]
    for i, (num, label, sub, color) in enumerate(steps):
        x = 0.52 + i * 1.80
        add_card(slide, x, 1.67, 1.56, 1.75, fill=LIGHT)
        add_text(slide, num, x + 0.20, 1.91, 0.30, 0.16, size=8.5, color=color, bold=True)
        add_shape(slide, MSO_SHAPE.OVAL, x + 1.16, 1.86, 0.18, 0.18, fill=color)
        add_text(slide, label, x + 0.20, 2.37, 1.15, 0.30, size=9.3, bold=True)
        add_text(slide, sub, x + 0.20, 2.88, 1.14, 0.28, size=8.8, color=GRAY)
    add_card(slide, 0.52, 3.72, 8.96, 1.13, fill=WHITE, line=LIGHT_2)
    add_text(slide, "TIME", 0.83, 4.03, 0.48, 0.15, size=7.5, color=GRAY, bold=True)
    segments = [
        (1.40, 4.00, 4.18, "TRAIN", GREEN),
        (5.58, 4.00, 1.48, "VALIDATION", ORANGE),
        (7.06, 4.00, 1.86, "TEST", PURPLE),
    ]
    for x, y, w, label, color in segments:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, 0.30, fill=color)
        add_text(slide, label, x, y + 0.08, w, 0.13, size=7.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(
        slide,
        "Never leak future information into the past.",
        5.26,
        4.51,
        3.66,
        0.17,
        size=8.5,
        color=RED,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )
    add_text(
        slide,
        "Compare accuracy, stability, cost and decision value.",
        0.72,
        5.10,
        8.58,
        0.20,
        size=10.5,
        color=DARK_GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    return slide


def slide_12(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[13])
    set_placeholder(slide.placeholders[10], "Thank you.", size=29, bold=False)
    set_placeholder(
        slide.placeholders[11],
        "Questions?\n\nStart with a strong baseline. Add complexity only when the evidence supports it.",
        size=13,
        color=GRAY,
    )
    return slide


def validate(prs: Presentation) -> None:
    assert len(prs.slides) == 12
    assert round(prs.slide_width / 914400, 3) == W
    assert round(prs.slide_height / 914400, 3) == H
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            assert shape.left >= 0, (slide_idx, shape.name, "left")
            assert shape.top >= 0, (slide_idx, shape.name, "top")
            assert shape.left + shape.width <= prs.slide_width + Inches(0.01), (
                slide_idx,
                shape.name,
                "right",
            )
            assert shape.top + shape.height <= prs.slide_height + Inches(0.01), (
                slide_idx,
                shape.name,
                "bottom",
            )


def main() -> None:
    prs = Presentation(TEMPLATE)
    clear_existing_slides(prs)
    prs.core_properties.title = "Machine Learning Models"
    prs.core_properties.subject = "Share Knowledge deck using the SiDi corporate template"
    prs.core_properties.author = "SiDi"

    slide_01(prs)
    slide_02(prs)
    slide_03(prs)
    slide_04(prs)
    slide_05(prs)
    slide_06(prs)
    slide_07(prs)
    slide_08(prs)
    slide_09(prs)
    slide_10(prs)
    slide_11(prs)
    slide_12(prs)

    validate(prs)
    prs.save(OUTPUT)
    print(f"Presentation created: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Size: {OUTPUT.stat().st_size / 1024:.1f} KiB")


if __name__ == "__main__":
    main()
