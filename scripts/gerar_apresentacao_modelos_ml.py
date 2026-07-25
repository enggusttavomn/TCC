#!/usr/bin/env python3
"""Gera uma apresentação 16:9 sobre modelos de Machine Learning.

Uso:
    PYTHONPATH=/tmp/tcc_pptx_lib python scripts/gerar_apresentacao_modelos_ml.py
"""

from __future__ import annotations

import math
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "apresentacao_modelos_machine_learning.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5

BG = "071421"
BG_ALT = "0A1B2B"
CARD = "10283C"
CARD_2 = "15334A"
WHITE = "F5F8FC"
MUTED = "A8BBCC"
GRID = "27465C"
CYAN = "2BD9C5"
ORANGE = "FF8B5C"
LIME = "C5F36A"
PURPLE = "9C8BFF"
BLUE = "4DA3FF"
PINK = "FF6FAE"
YELLOW = "FFD166"
RED = "FF6174"
FONT = "Lato"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color.replace("#", ""))


def add_shape(
    slide,
    shape_type,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str | None = None,
    line: str | None = None,
    line_width: float = 1.0,
    radius: bool = False,
):
    if radius:
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color: str = WHITE,
    bold: bool = False,
    font: str = FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.0,
    tracking: float | None = None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    p.line_spacing = 1.0
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    if tracking is not None:
        # PowerPoint stores tracking in 1/1000 pt.
        run._r.get_or_add_rPr().set("spc", str(int(tracking * 1000)))
    return box


def add_line(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    *,
    color: str = GRID,
    width: float = 1.5,
    dash: str | None = None,
):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    if dash:
        line.line.dash_style = dash
    return line


def add_chevron(slide, x: float, y: float, color: str, size: float = 0.13):
    shape = add_shape(
        slide,
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        x,
        y,
        size,
        size,
        fill=color,
    )
    shape.rotation = 90
    return shape


def add_pill(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    *,
    fill: str = CARD_2,
    color: str = WHITE,
    size: float = 11,
    line: str | None = None,
):
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x,
        y,
        w,
        0.34,
        fill=fill,
        line=line,
        radius=True,
    )
    add_text(
        slide,
        text,
        x,
        y + 0.01,
        w,
        0.29,
        size=size,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_header(slide, kicker: str, title: str, number: int):
    add_text(
        slide,
        kicker.upper(),
        0.62,
        0.35,
        5.0,
        0.28,
        size=10.5,
        color=CYAN,
        bold=True,
        tracking=1.2,
    )
    add_text(slide, title, 0.62, 0.68, 11.9, 0.65, size=28, bold=True)
    add_text(
        slide,
        f"{number:02d}",
        12.15,
        0.39,
        0.55,
        0.3,
        size=11,
        color=MUTED,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )
    add_line(slide, 0.62, 1.38, 12.71, 1.38, color=GRID, width=1)


def add_footer(slide, section: str):
    add_text(
        slide,
        "SHARE KNOWLEDGE  •  MODELOS DE MACHINE LEARNING",
        0.62,
        7.1,
        5.7,
        0.18,
        size=8.5,
        color="6F879A",
        bold=True,
        tracking=0.5,
    )
    add_text(
        slide,
        section.upper(),
        9.2,
        7.1,
        3.5,
        0.18,
        size=8.5,
        color="6F879A",
        bold=True,
        align=PP_ALIGN.RIGHT,
        tracking=0.5,
    )


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = CARD,
    accent: str | None = None,
):
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x,
        y,
        w,
        h,
        fill=fill,
        radius=True,
    )
    if accent:
        add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x,
            y,
            0.08,
            h,
            fill=accent,
            radius=True,
        )


def add_bullet(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    *,
    color: str = WHITE,
    dot: str = CYAN,
    size: float = 14,
):
    add_shape(slide, MSO_SHAPE.OVAL, x, y + 0.11, 0.08, 0.08, fill=dot)
    add_text(slide, text, x + 0.18, y, w - 0.18, 0.42, size=size, color=color)


def add_model_tag(slide, name: str, x: float, y: float, color: str):
    add_pill(slide, name, x, y, 1.45, fill=CARD_2, color=color, size=11, line=color)


def new_slide(prs, bg: str = BG):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = rgb(bg)
    return slide


def slide_01(prs):
    slide = new_slide(prs)
    # Decorative orbit.
    add_shape(slide, MSO_SHAPE.OVAL, 9.62, -0.82, 4.6, 4.6, line=GRID, line_width=1)
    add_shape(slide, MSO_SHAPE.OVAL, 10.3, -0.14, 3.24, 3.24, line=GRID, line_width=1)
    for x, y, c, s in [
        (11.72, 0.34, CYAN, 0.25),
        (10.42, 1.16, ORANGE, 0.18),
        (12.34, 2.35, LIME, 0.16),
        (9.7, 2.02, PURPLE, 0.14),
    ]:
        add_shape(slide, MSO_SHAPE.OVAL, x, y, s, s, fill=c)

    add_text(
        slide,
        "SHARE KNOWLEDGE",
        0.72,
        0.72,
        3.2,
        0.3,
        size=11,
        color=CYAN,
        bold=True,
        tracking=1.5,
    )
    add_text(
        slide,
        "Modelos de\nMachine Learning",
        0.72,
        1.36,
        8.1,
        1.85,
        size=42,
        color=WHITE,
        bold=True,
    )
    add_text(
        slide,
        "Sete formas de aprender padrões — e o que muda na hora de prever.",
        0.76,
        3.42,
        7.4,
        0.56,
        size=18,
        color=MUTED,
    )
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        0.72,
        4.33,
        11.92,
        1.46,
        fill=BG_ALT,
        line=GRID,
        radius=True,
    )
    models = [
        ("XGBoost", ORANGE),
        ("MLP", BLUE),
        ("RNN", CYAN),
        ("LSTM", LIME),
        ("DilatedRNN", PURPLE),
        ("DeepAR", PINK),
        ("DeepNPTS*", YELLOW),
    ]
    for idx, (name, color) in enumerate(models):
        x = 0.98 + idx * 1.64
        add_shape(slide, MSO_SHAPE.OVAL, x + 0.56, 4.65, 0.20, 0.20, fill=color)
        add_text(
            slide,
            name,
            x,
            5.05,
            1.34,
            0.30,
            size=11.5,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        if idx < len(models) - 1:
            add_line(slide, x + 1.22, 4.75, x + 1.72, 4.75, color=GRID, width=1)
    add_text(
        slide,
        "* aproximação experimental inspirada no DeepNPTS",
        0.76,
        6.16,
        5.5,
        0.26,
        size=10.5,
        color=MUTED,
    )
    add_text(
        slide,
        "De padrões tabulares a previsões probabilísticas",
        7.3,
        6.12,
        5.3,
        0.30,
        size=10.5,
        color=CYAN,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )
    return slide


def slide_02(prs):
    slide = new_slide(prs)
    add_header(slide, "Visão geral", "A mesma pergunta, sete vieses de aprendizagem", 2)
    add_text(
        slide,
        "O modelo muda o que é fácil perceber nos dados.",
        0.62,
        1.62,
        6.2,
        0.38,
        size=16,
        color=MUTED,
    )

    axes = [
        (
            "01",
            "Como os dados entram?",
            "Tabela de atributos",
            "Sequência ordenada",
            ORANGE,
        ),
        (
            "02",
            "Como o passado é lembrado?",
            "Sem estado interno",
            "Memória curta, longa ou multiescala",
            CYAN,
        ),
        (
            "03",
            "O que sai do modelo?",
            "Um valor pontual",
            "Uma distribuição de futuros",
            PURPLE,
        ),
    ]
    for i, (num, title, left, right, accent) in enumerate(axes):
        y = 2.18 + i * 1.42
        add_card(slide, 0.62, y, 12.08, 1.08, fill=CARD, accent=accent)
        add_text(slide, num, 0.95, y + 0.20, 0.55, 0.34, size=14, color=accent, bold=True)
        add_text(slide, title, 1.55, y + 0.17, 3.2, 0.34, size=16, bold=True)
        add_pill(slide, left, 5.05, y + 0.20, 2.85, fill=BG_ALT, color=WHITE, size=11)
        add_line(slide, 8.1, y + 0.52, 8.55, y + 0.52, color=accent, width=2)
        add_chevron(slide, 8.51, y + 0.455, accent)
        add_pill(slide, right, 8.72, y + 0.20, 3.55, fill=BG_ALT, color=WHITE, size=11)

    add_text(
        slide,
        "A escolha do modelo é uma hipótese sobre a estrutura do problema.",
        0.86,
        6.55,
        11.6,
        0.34,
        size=16,
        color=CYAN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, "Mapa mental")
    return slide


def draw_trees(slide, x: float, y: float):
    colors = [ORANGE, YELLOW, PINK]
    for k, color in enumerate(colors):
        ox = x + k * 1.15
        add_line(slide, ox + 0.44, y + 0.25, ox + 0.20, y + 0.67, color=GRID, width=1.3)
        add_line(slide, ox + 0.44, y + 0.25, ox + 0.68, y + 0.67, color=GRID, width=1.3)
        add_line(slide, ox + 0.20, y + 0.67, ox + 0.08, y + 1.02, color=GRID, width=1.3)
        add_line(slide, ox + 0.20, y + 0.67, ox + 0.33, y + 1.02, color=GRID, width=1.3)
        add_shape(slide, MSO_SHAPE.OVAL, ox + 0.32, y + 0.12, 0.24, 0.24, fill=color)
        add_shape(slide, MSO_SHAPE.OVAL, ox + 0.09, y + 0.56, 0.22, 0.22, fill=CARD_2, line=color)
        add_shape(slide, MSO_SHAPE.OVAL, ox + 0.57, y + 0.56, 0.22, 0.22, fill=CARD_2, line=color)
        add_shape(slide, MSO_SHAPE.OVAL, ox, y + 0.96, 0.16, 0.16, fill=color)
        add_shape(slide, MSO_SHAPE.OVAL, ox + 0.27, y + 0.96, 0.16, 0.16, fill=color)


def draw_mlp(slide, x: float, y: float):
    layer_x = [x, x + 1.0, x + 2.0]
    layer_count = [4, 5, 2]
    layer_colors = [BLUE, CYAN, LIME]
    coords = []
    for lx, count, color in zip(layer_x, layer_count, layer_colors):
        ys = [y + j * (1.22 / (count - 1)) for j in range(count)] if count > 1 else [y]
        coords.append([(lx, yy) for yy in ys])
        for _, yy in coords[-1]:
            add_shape(slide, MSO_SHAPE.OVAL, lx, yy, 0.20, 0.20, fill=color)
    for left, right in zip(coords[:-1], coords[1:]):
        for x1, y1 in left:
            for x2, y2 in right:
                add_line(slide, x1 + 0.20, y1 + 0.10, x2, y2 + 0.10, color=GRID, width=0.5)


def slide_03(prs):
    slide = new_slide(prs)
    add_header(slide, "Modelos tabulares", "Dois caminhos quando a entrada já virou atributos", 3)

    add_card(slide, 0.62, 1.65, 5.85, 4.92, fill=CARD)
    add_model_tag(slide, "XGBoost", 0.95, 1.95, ORANGE)
    add_text(
        slide,
        "Muitas árvores pequenas\ncorrigem os erros anteriores.",
        0.95,
        2.48,
        3.65,
        0.75,
        size=18,
        bold=True,
    )
    draw_trees(slide, 1.16, 3.43)
    add_text(slide, "Σ", 4.66, 3.68, 0.45, 0.45, size=27, color=ORANGE, bold=True)
    add_line(slide, 4.06, 3.94, 4.60, 3.94, color=ORANGE, width=2)
    add_line(slide, 5.05, 3.94, 5.48, 3.94, color=ORANGE, width=2)
    add_shape(slide, MSO_SHAPE.OVAL, 5.46, 3.78, 0.30, 0.30, fill=ORANGE)
    add_bullet(slide, "Forte em dados tabulares e interações.", 0.96, 5.05, 4.9, dot=ORANGE)
    add_bullet(slide, "Pouca preparação; boa referência inicial.", 0.96, 5.48, 4.9, dot=ORANGE)
    add_text(
        slide,
        "ATENÇÃO  •  a sequência precisa virar lags, médias e calendário.",
        0.96,
        6.05,
        4.9,
        0.28,
        size=9.5,
        color=MUTED,
        bold=True,
    )

    add_card(slide, 6.83, 1.65, 5.87, 4.92, fill=CARD)
    add_model_tag(slide, "MLP", 7.16, 1.95, BLUE)
    add_text(
        slide,
        "Camadas densas transformam\natributos em relações não lineares.",
        7.16,
        2.48,
        4.25,
        0.75,
        size=18,
        bold=True,
    )
    draw_mlp(slide, 7.55, 3.50)
    add_line(slide, 9.82, 4.21, 11.56, 4.21, color=LIME, width=2)
    add_chevron(slide, 11.52, 4.145, LIME)
    add_shape(slide, MSO_SHAPE.OVAL, 11.78, 4.05, 0.32, 0.32, fill=LIME)
    add_bullet(slide, "Flexível para relações suaves e complexas.", 7.16, 5.05, 4.9, dot=BLUE)
    add_bullet(slide, "Exige escala, regularização e ajuste.", 7.16, 5.48, 4.9, dot=BLUE)
    add_text(
        slide,
        "ATENÇÃO  •  sem recorrência, não mantém estado temporal por si só.",
        7.16,
        6.05,
        4.9,
        0.28,
        size=9.5,
        color=MUTED,
        bold=True,
    )
    add_footer(slide, "XGBoost + MLP")
    return slide


def slide_04(prs):
    slide = new_slide(prs)
    add_header(slide, "Sequências", "RNN: o estado carrega contexto", 4)
    add_text(
        slide,
        "Cada passo combina a observação atual com um resumo do passado.",
        0.62,
        1.62,
        7.4,
        0.36,
        size=16,
        color=MUTED,
    )

    y = 3.12
    for i in range(5):
        x = 0.92 + i * 2.15
        add_pill(slide, f"xₜ₋{4-i}" if i < 4 else "xₜ", x + 0.34, 2.12, 0.72, fill=BG_ALT, color=CYAN, size=13)
        add_line(slide, x + 0.70, 2.47, x + 0.70, y - 0.08, color=CYAN, width=1.7)
        add_chevron(slide, x + 0.635, y - 0.14, CYAN)
        add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x,
            y,
            1.42,
            1.05,
            fill=CARD_2,
            line=CYAN if i == 4 else GRID,
            line_width=1.5,
            radius=True,
        )
        add_text(
            slide,
            "RNN",
            x,
            y + 0.20,
            1.42,
            0.28,
            size=15,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            f"h{i+1}",
            x,
            y + 0.60,
            1.42,
            0.22,
            size=10,
            color=MUTED,
            align=PP_ALIGN.CENTER,
        )
        if i < 4:
            add_line(slide, x + 1.42, y + 0.52, x + 2.05, y + 0.52, color=ORANGE, width=2)
            add_chevron(slide, x + 2.00, y + 0.455, ORANGE)

    add_line(slide, 10.24, y + 1.06, 10.24, 5.02, color=LIME, width=2)
    add_chevron(slide, 10.175, 4.97, LIME)
    add_pill(slide, "previsão ŷₜ₊₁", 9.40, 5.15, 1.70, fill=LIME, color=BG, size=12)

    add_card(slide, 0.92, 5.15, 7.35, 1.18, fill=BG_ALT, accent=ORANGE)
    add_text(slide, "FORÇA", 1.20, 5.39, 0.85, 0.22, size=10, color=ORANGE, bold=True)
    add_text(
        slide,
        "Ordem e dependências locais entram diretamente na arquitetura.",
        2.10,
        5.29,
        5.78,
        0.45,
        size=14,
        bold=True,
    )
    add_text(slide, "LIMITE", 1.20, 5.87, 0.85, 0.22, size=10, color=RED, bold=True)
    add_text(
        slide,
        "Sinais distantes podem enfraquecer durante o treinamento.",
        2.10,
        5.78,
        5.78,
        0.40,
        size=13,
        color=MUTED,
    )
    add_footer(slide, "RNN")
    return slide


def slide_05(prs):
    slide = new_slide(prs)
    add_header(slide, "Memória seletiva", "LSTM: lembrar também significa esquecer", 5)
    add_text(
        slide,
        "A célula controla quais informações entram, permanecem e saem.",
        0.62,
        1.62,
        7.6,
        0.36,
        size=16,
        color=MUTED,
    )

    add_card(slide, 0.72, 2.18, 8.28, 3.84, fill=CARD)
    # Cell-state highway.
    add_text(slide, "estado da memória  cₜ₋₁", 1.10, 2.47, 2.2, 0.28, size=11, color=MUTED, bold=True)
    add_line(slide, 1.12, 3.02, 8.52, 3.02, color=LIME, width=4)
    add_chevron(slide, 8.43, 2.93, LIME, size=0.18)
    add_text(slide, "cₜ", 8.22, 2.47, 0.4, 0.28, size=12, color=LIME, bold=True, align=PP_ALIGN.CENTER)

    gates = [
        (1.28, "ESQUECER", "o que sai", ORANGE, "×"),
        (3.53, "ENTRADA", "o que entra", CYAN, "+"),
        (5.78, "SAÍDA", "o que vira hₜ", PURPLE, "σ"),
    ]
    for x, label, sub, color, symbol in gates:
        add_shape(slide, MSO_SHAPE.OVAL, x, 3.66, 0.62, 0.62, fill=BG_ALT, line=color, line_width=2)
        add_text(
            slide,
            symbol,
            x,
            3.79,
            0.62,
            0.25,
            size=17,
            color=color,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_line(slide, x + 0.31, 3.66, x + 0.31, 3.06, color=color, width=1.8)
        add_text(slide, label, x - 0.25, 4.53, 1.12, 0.23, size=9, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, sub, x - 0.45, 4.82, 1.52, 0.28, size=11, color=MUTED, align=PP_ALIGN.CENTER)

    add_pill(slide, "xₜ + hₜ₋₁", 3.70, 5.36, 1.72, fill=BG_ALT, color=WHITE, size=12, line=GRID)
    for x in [1.59, 3.84, 6.09]:
        add_line(slide, 4.56, 5.36, x, 4.30, color=GRID, width=1)

    add_card(slide, 9.37, 2.18, 3.34, 3.84, fill=BG_ALT, accent=LIME)
    add_text(slide, "POR QUE FUNCIONA?", 9.73, 2.55, 2.5, 0.25, size=10, color=LIME, bold=True)
    add_text(
        slide,
        "O caminho da memória reduz a perda de sinais relevantes em horizontes maiores.",
        9.73,
        3.04,
        2.43,
        1.12,
        size=18,
        bold=True,
    )
    add_text(
        slide,
        "Trade-off",
        9.73,
        4.62,
        1.2,
        0.25,
        size=11,
        color=ORANGE,
        bold=True,
    )
    add_text(
        slide,
        "Mais parâmetros e treinamento mais cuidadoso que uma RNN simples.",
        9.73,
        5.00,
        2.43,
        0.74,
        size=13,
        color=MUTED,
    )
    add_footer(slide, "LSTM")
    return slide


def slide_06(prs):
    slide = new_slide(prs)
    add_header(slide, "Múltiplas escalas", "DilatedRNN: conexões que saltam no tempo", 6)
    add_text(
        slide,
        "O estado consulta diretamente hₜ₋ₛ; diferentes saltos capturam ritmos diferentes.",
        0.62,
        1.62,
        10.4,
        0.36,
        size=16,
        color=MUTED,
    )

    # Timeline nodes.
    xs = [0.98 + i * 1.04 for i in range(11)]
    base_y = 5.30
    for i, x in enumerate(xs):
        color = PURPLE if i == 10 else CARD_2
        outline = PURPLE if i in [2, 6, 8, 9, 10] else GRID
        add_shape(slide, MSO_SHAPE.OVAL, x, base_y, 0.36, 0.36, fill=color, line=outline, line_width=1.2)
        add_text(slide, str(i), x - 0.03, base_y + 0.47, 0.42, 0.20, size=9, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, "tempo", 11.98, 5.72, 0.6, 0.22, size=10, color=MUTED)

    # Arcs approximated by elbow connectors.
    levels = [
        (1, 4.73, CYAN, "s = 1  •  curto prazo"),
        (2, 3.90, ORANGE, "s = 2  •  padrão intermediário"),
        (4, 2.93, PURPLE, "s = 4  •  padrão mais amplo"),
    ]
    target_idx = 10
    for step, y, color, label in levels:
        source_idx = target_idx - step
        x1 = xs[source_idx] + 0.18
        x2 = xs[target_idx] + 0.18
        add_line(slide, x1, base_y, x1, y, color=color, width=2)
        add_line(slide, x1, y, x2, y, color=color, width=2)
        add_line(slide, x2, y, x2, base_y, color=color, width=2)
        add_chevron(slide, x2 - 0.065, base_y - 0.14, color)
        add_pill(slide, label, 0.98, y - 0.18, 2.55, fill=BG_ALT, color=color, size=10, line=color)

    add_card(slide, 8.48, 1.86, 4.22, 1.55, fill=BG_ALT, accent=PURPLE)
    add_text(slide, "RECORRÊNCIA DILATADA", 8.84, 2.15, 2.9, 0.23, size=10, color=PURPLE, bold=True)
    add_text(
        slide,
        "hₜ = φ(Wₓxₜ + Wₕhₜ₋ₛ + b)",
        8.84,
        2.61,
        3.34,
        0.40,
        size=18,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    add_card(slide, 0.98, 6.20, 11.20, 0.60, fill=CARD)
    add_text(
        slide,
        "NA IMPLEMENTAÇÃO  •  sequência real entra na recorrência; covariáveis auxiliares são combinadas depois do codificador.",
        1.20,
        6.39,
        10.76,
        0.22,
        size=10.5,
        color=MUTED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, "DilatedRNN")
    return slide


def slide_07(prs):
    slide = new_slide(prs)
    add_header(slide, "Previsão probabilística", "DeepAR: não prever um número, mas futuros possíveis", 7)
    add_text(
        slide,
        "Um modelo global aprende padrões compartilhados e produz uma distribuição por passo.",
        0.62,
        1.62,
        10.8,
        0.36,
        size=16,
        color=MUTED,
    )

    add_card(slide, 0.72, 2.15, 3.48, 3.98, fill=CARD)
    add_text(slide, "MUITAS SÉRIES", 1.08, 2.48, 2.2, 0.24, size=10, color=CYAN, bold=True)
    series_colors = [CYAN, ORANGE, PURPLE, PINK]
    for row, color in enumerate(series_colors):
        yy = 3.08 + row * 0.58
        points = []
        for i in range(8):
            xx = 1.07 + i * 0.34
            val = 0.11 * math.sin(i * 1.3 + row) + 0.025 * i
            points.append((xx, yy - val))
        for p1, p2 in zip(points[:-1], points[1:]):
            add_line(slide, *p1, *p2, color=color, width=1.8)
    add_pill(slide, "parâmetros compartilhados", 1.08, 5.52, 2.74, fill=BG_ALT, color=WHITE, size=10, line=GRID)

    add_line(slide, 4.21, 4.12, 4.93, 4.12, color=CYAN, width=2.2)
    add_chevron(slide, 4.87, 4.055, CYAN)
    add_shape(slide, MSO_SHAPE.OVAL, 5.08, 3.42, 1.38, 1.38, fill=CARD_2, line=CYAN, line_width=2)
    add_text(slide, "RNN", 5.08, 3.69, 1.38, 0.30, size=18, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "global", 5.08, 4.06, 1.38, 0.23, size=10, color=MUTED, align=PP_ALIGN.CENTER)
    add_line(slide, 6.46, 4.12, 7.18, 4.12, color=CYAN, width=2.2)
    add_chevron(slide, 7.12, 4.055, CYAN)

    add_card(slide, 7.34, 2.15, 5.36, 3.98, fill=BG_ALT)
    add_text(slide, "DISTRIBUIÇÃO DO FUTURO", 7.72, 2.48, 3.2, 0.24, size=10, color=PINK, bold=True)
    chart_x0, chart_y0 = 7.75, 5.55
    add_line(slide, chart_x0, chart_y0, 12.28, chart_y0, color=GRID, width=1)
    add_line(slide, chart_x0, 2.97, chart_x0, chart_y0, color=GRID, width=1)
    historical = []
    for i in range(9):
        xx = chart_x0 + i * 0.28
        yy = 4.92 - 0.12 * i - 0.20 * math.sin(i * 0.85)
        historical.append((xx, yy))
    for p1, p2 in zip(historical[:-1], historical[1:]):
        add_line(slide, *p1, *p2, color=WHITE, width=2)
    split_x = historical[-1][0]
    add_line(slide, split_x, 2.97, split_x, chart_y0, color=MUTED, width=1)
    add_text(slide, "agora", split_x - 0.27, 5.68, 0.58, 0.18, size=8.5, color=MUTED, align=PP_ALIGN.CENTER)
    median = historical[-1][1]
    future_xs = [split_x + i * 0.31 for i in range(1, 7)]
    upper, lower, med = [], [], []
    for i, xx in enumerate(future_xs, 1):
        center = median - 0.10 * i + 0.05 * math.sin(i)
        spread = 0.08 + 0.055 * i
        upper.append((xx, center - spread))
        lower.append((xx, center + spread))
        med.append((xx, center))
        add_line(slide, xx, center - spread, xx, center + spread, color="503C61", width=6)
    for arr, color, width in [(upper, PINK, 1), (lower, PINK, 1), (med, YELLOW, 2.5)]:
        start = historical[-1]
        for p1, p2 in zip([start] + arr[:-1], arr):
            add_line(slide, *p1, *p2, color=color, width=width)
    add_pill(slide, "P10", 10.92, 3.08, 0.62, fill=CARD_2, color=PINK, size=9)
    add_pill(slide, "P50", 11.58, 3.08, 0.62, fill=CARD_2, color=YELLOW, size=9)
    add_pill(slide, "P90", 10.92, 3.48, 0.62, fill=CARD_2, color=PINK, size=9)

    add_text(
        slide,
        "Útil quando o custo da incerteza importa: estoque, capacidade, risco e energia.",
        0.92,
        6.42,
        11.5,
        0.33,
        size=14,
        color=CYAN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, "DeepAR")
    return slide


def slide_08(prs):
    slide = new_slide(prs)
    add_header(
        slide,
        "Abordagem experimental",
        "DeepNPTS-inspired: o histórico vira um conjunto de candidatos",
        8,
    )
    add_text(
        slide,
        "Em vez de gerar qualquer valor, a previsão reutiliza observações passadas com pesos aprendidos.",
        0.62,
        1.62,
        11.1,
        0.36,
        size=16,
        color=MUTED,
    )

    stages = [
        (0.72, "1", "CONTEXTO", "padrão recente", CYAN),
        (3.50, "2", "SCORE", "similaridade + recência", ORANGE),
        (6.42, "3", "PESOS", "softmax sobre o histórico", PURPLE),
        (9.34, "4", "PREVISÃO", "amostra valores observados", YELLOW),
    ]
    for i, (x, num, label, sub, color) in enumerate(stages):
        add_card(slide, x, 2.42, 2.28, 3.28, fill=CARD)
        add_shape(slide, MSO_SHAPE.OVAL, x + 0.18, 2.61, 0.36, 0.36, fill=color)
        add_text(slide, num, x + 0.18, 2.68, 0.36, 0.18, size=10, color=BG, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, label, x + 0.68, 2.66, 1.35, 0.22, size=10, color=color, bold=True)
        add_text(slide, sub, x + 0.20, 5.13, 1.88, 0.34, size=11, color=MUTED, align=PP_ALIGN.CENTER)
        if i < 3:
            add_line(slide, x + 2.28, 4.02, x + 2.70, 4.02, color=color, width=2)
            add_chevron(slide, x + 2.65, 3.955, color)

    # Context mini-series.
    pts = []
    for i in range(8):
        xx = 1.03 + i * 0.21
        yy = 4.44 - 0.27 * math.sin(i * 0.95) - 0.035 * i
        pts.append((xx, yy))
    for p1, p2 in zip(pts[:-1], pts[1:]):
        add_line(slide, *p1, *p2, color=CYAN, width=2)
    # Scoring bars.
    for i, width in enumerate([0.38, 0.82, 1.28, 0.60]):
        add_shape(slide, MSO_SHAPE.RECTANGLE, 3.86, 3.42 + i * 0.34, width, 0.16, fill=ORANGE)
    # Weighted history dots.
    weights = [0.10, 0.42, 0.16, 0.72, 0.30, 0.92]
    for i, val in enumerate(weights):
        size = 0.12 + val * 0.22
        add_shape(slide, MSO_SHAPE.OVAL, 6.78 + i * 0.28, 4.19 - size / 2, size, size, fill=PURPLE)
    add_text(slide, "Σ pᵢ = 1", 7.00, 3.38, 1.25, 0.32, size=18, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
    # Candidate values / resulting discrete distribution.
    values = [(9.76, 4.58, 0.19), (10.22, 3.82, 0.30), (10.74, 4.34, 0.24), (11.24, 3.52, 0.38)]
    for x, y, size in values:
        add_line(slide, x + size / 2, 4.86, x + size / 2, y + size / 2, color=GRID, width=1)
        add_shape(slide, MSO_SHAPE.OVAL, x, y, size, size, fill=YELLOW)
    add_line(slide, 9.70, 4.86, 11.88, 4.86, color=GRID, width=1.2)

    add_card(slide, 0.72, 6.02, 11.86, 0.70, fill=BG_ALT, accent=YELLOW)
    add_text(
        slide,
        "IMPORTANTE  •  aqui tratamos uma aproximação experimental inspirada no princípio não paramétrico — não uma reprodução integral do DeepNPTS oficial.",
        1.02,
        6.24,
        11.14,
        0.27,
        size=10.5,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, "DeepNPTS-inspired")
    return slide


def slide_09(prs):
    slide = new_slide(prs)
    add_header(slide, "Comparação", "Os sete modelos em uma tela", 9)

    x0 = 0.62
    widths = [1.55, 2.22, 2.18, 1.58, 2.58, 1.94]
    headers = ["MODELO", "ENTRADA", "MEMÓRIA / ESCALA", "SAÍDA", "PONTO FORTE", "CUIDADO"]
    total_w = sum(widths)
    cur = x0
    for h, w in zip(headers, widths):
        add_shape(slide, MSO_SHAPE.RECTANGLE, cur, 1.67, w, 0.50, fill=CARD_2)
        add_text(slide, h, cur + 0.12, 1.82, w - 0.24, 0.18, size=8.5, color=MUTED, bold=True)
        cur += w

    rows = [
        ("XGBoost", "atributos", "lags explícitos", "pontual", "tabular robusto", "engenharia temporal", ORANGE),
        ("MLP", "atributos", "lags explícitos", "pontual", "não linearidade", "escala / ajuste", BLUE),
        ("RNN", "sequência", "curta", "pontual", "ordem temporal", "sinais distantes", CYAN),
        ("LSTM", "sequência", "longa seletiva", "pontual", "dependência longa", "mais parâmetros", LIME),
        ("DilatedRNN", "sequência", "multiescala", "pontual", "ritmos distintos", "saltos adequados", PURPLE),
        ("DeepAR", "muitas séries", "autoregressiva", "distribuição", "incerteza", "hipótese da dist.", PINK),
        ("DeepNPTS*", "histórico", "similaridade", "distribuição", "valores plausíveis", "limite ao histórico", YELLOW),
    ]
    row_h = 0.60
    for r, row in enumerate(rows):
        y = 2.17 + r * row_h
        fill = BG_ALT if r % 2 else CARD
        cur = x0
        for c, (value, w) in enumerate(zip(row[:-1], widths)):
            add_shape(slide, MSO_SHAPE.RECTANGLE, cur, y, w, row_h, fill=fill)
            color = row[-1] if c == 0 else (WHITE if c in [1, 2, 3] else MUTED)
            add_text(
                slide,
                value,
                cur + 0.12,
                y + 0.16,
                w - 0.24,
                0.24,
                size=10 if c else 10.5,
                color=color,
                bold=(c == 0),
            )
            cur += w
        add_shape(slide, MSO_SHAPE.RECTANGLE, x0, y, 0.06, row_h, fill=row[-1])

    add_text(
        slide,
        "* aproximação experimental inspirada no DeepNPTS",
        0.64,
        6.54,
        4.8,
        0.22,
        size=9.5,
        color=MUTED,
    )
    add_pill(slide, "pontual = um melhor palpite", 7.20, 6.44, 2.35, fill=CARD, color=WHITE, size=9.5)
    add_pill(slide, "distribuição = cenários + risco", 9.72, 6.44, 2.72, fill=CARD, color=WHITE, size=9.5)
    add_footer(slide, "Comparativo")
    return slide


def slide_10(prs):
    slide = new_slide(prs)
    add_header(slide, "Guia de escolha", "Qual modelo testar primeiro?", 10)
    add_text(
        slide,
        "Use o requisito do problema para reduzir o espaço de busca — depois valide com dados.",
        0.62,
        1.62,
        10.7,
        0.36,
        size=16,
        color=MUTED,
    )

    add_card(slide, 0.72, 2.35, 3.00, 1.02, fill=CARD_2, accent=CYAN)
    add_text(slide, "Precisa quantificar incerteza?", 1.08, 2.68, 2.28, 0.30, size=15, bold=True, align=PP_ALIGN.CENTER)
    add_line(slide, 3.72, 2.86, 4.42, 2.86, color=CYAN, width=2)
    add_chevron(slide, 4.36, 2.795, CYAN)

    add_pill(slide, "SIM", 4.52, 2.38, 0.72, fill=CYAN, color=BG, size=10)
    add_card(slide, 5.44, 2.18, 3.16, 1.38, fill=CARD)
    add_text(slide, "DeepAR", 5.74, 2.48, 1.06, 0.28, size=17, color=PINK, bold=True)
    add_text(slide, "distribuição paramétrica", 5.74, 2.87, 2.45, 0.25, size=11, color=MUTED)
    add_card(slide, 8.84, 2.18, 3.82, 1.38, fill=CARD)
    add_text(slide, "DeepNPTS-inspired", 9.14, 2.48, 2.24, 0.28, size=17, color=YELLOW, bold=True)
    add_text(slide, "massa sobre valores históricos", 9.14, 2.87, 2.86, 0.25, size=11, color=MUTED)

    add_line(slide, 2.22, 3.37, 2.22, 4.18, color=ORANGE, width=2)
    add_chevron(slide, 2.155, 4.11, ORANGE)
    add_pill(slide, "NÃO", 2.42, 3.69, 0.72, fill=ORANGE, color=BG, size=10)
    add_card(slide, 0.72, 4.32, 3.00, 1.02, fill=CARD_2, accent=ORANGE)
    add_text(slide, "A ordem temporal é central?", 1.08, 4.65, 2.28, 0.30, size=15, bold=True, align=PP_ALIGN.CENTER)

    add_line(slide, 3.72, 4.83, 4.42, 4.83, color=ORANGE, width=2)
    add_chevron(slide, 4.36, 4.765, ORANGE)
    add_pill(slide, "NÃO", 4.52, 4.35, 0.72, fill=ORANGE, color=BG, size=10)
    add_card(slide, 5.44, 4.15, 3.16, 1.38, fill=CARD)
    add_text(slide, "XGBoost", 5.74, 4.45, 1.35, 0.28, size=17, color=ORANGE, bold=True)
    add_text(slide, "baseline tabular forte", 5.74, 4.84, 2.45, 0.25, size=11, color=MUTED)
    add_card(slide, 8.84, 4.15, 3.82, 1.38, fill=CARD)
    add_text(slide, "MLP", 9.14, 4.45, 0.80, 0.28, size=17, color=BLUE, bold=True)
    add_text(slide, "quando há dados e não linearidade", 9.14, 4.84, 2.92, 0.25, size=11, color=MUTED)

    add_line(slide, 2.22, 5.34, 2.22, 6.10, color=LIME, width=2)
    add_chevron(slide, 2.155, 6.03, LIME)
    add_pill(slide, "SIM", 2.42, 5.66, 0.72, fill=LIME, color=BG, size=10)
    add_card(slide, 0.72, 6.15, 11.94, 0.63, fill=BG_ALT)
    add_text(slide, "RNN  →  baseline sequencial", 1.03, 6.35, 2.92, 0.22, size=11, color=CYAN, bold=True)
    add_text(slide, "LSTM  →  dependências longas", 4.46, 6.35, 3.00, 0.22, size=11, color=LIME, bold=True)
    add_text(slide, "DilatedRNN  →  múltiplas escalas", 8.00, 6.35, 3.68, 0.22, size=11, color=PURPLE, bold=True)
    add_footer(slide, "Decisão")
    return slide


def slide_11(prs):
    slide = new_slide(prs)
    add_header(slide, "Avaliação", "Um experimento justo vale mais que uma arquitetura sofisticada", 11)
    add_text(
        slide,
        "Fixe o protocolo; varie apenas o modelo.",
        0.62,
        1.62,
        7.4,
        0.36,
        size=16,
        color=MUTED,
    )

    steps = [
        ("01", "CORTE TEMPORAL", "treino antes, teste depois", CYAN),
        ("02", "MESMAS ENTRADAS", "lags e covariáveis comparáveis", ORANGE),
        ("03", "MESMO ORÇAMENTO", "busca e tempo controlados", PURPLE),
        ("04", "MÉTRICAS CERTAS", "erro + calibração + custo", LIME),
        ("05", "REPETIÇÕES", "sementes e intervalos", PINK),
    ]
    for i, (num, title, sub, color) in enumerate(steps):
        x = 0.70 + i * 2.48
        add_card(slide, x, 2.30, 2.18, 2.04, fill=CARD)
        add_text(slide, num, x + 0.20, 2.55, 0.42, 0.24, size=12, color=color, bold=True)
        add_shape(slide, MSO_SHAPE.OVAL, x + 1.62, 2.49, 0.24, 0.24, fill=color)
        add_text(slide, title, x + 0.20, 3.05, 1.76, 0.30, size=11.5, color=WHITE, bold=True)
        add_text(slide, sub, x + 0.20, 3.48, 1.76, 0.52, size=11.5, color=MUTED)

    # Temporal split visual.
    add_card(slide, 0.70, 4.72, 12.02, 1.42, fill=BG_ALT)
    add_text(slide, "TEMPO", 1.03, 5.00, 0.62, 0.22, size=9, color=MUTED, bold=True)
    segments = [
        (1.77, 4.98, 5.62, 0.38, "TREINO", CYAN),
        (7.39, 4.98, 2.01, 0.38, "VALIDAÇÃO", ORANGE),
        (9.40, 4.98, 2.64, 0.38, "TESTE", PURPLE),
    ]
    for x, y, w, h, label, color in segments:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=color, radius=True)
        add_text(slide, label, x, y + 0.09, w, 0.17, size=8.5, color=BG, bold=True, align=PP_ALIGN.CENTER)
    add_line(slide, 1.77, 5.67, 12.04, 5.67, color=GRID, width=1)
    for i in range(12):
        xx = 1.77 + i * 0.93
        add_line(slide, xx, 5.59, xx, 5.75, color=GRID, width=1)
    add_text(
        slide,
        "Sem embaralhar o futuro para dentro do passado.",
        7.28,
        5.82,
        4.76,
        0.22,
        size=10.5,
        color=RED,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )

    add_text(
        slide,
        "Só então compare acurácia, estabilidade, custo e utilidade para a decisão.",
        0.80,
        6.46,
        11.7,
        0.28,
        size=14,
        color=CYAN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, "Protocolo")
    return slide


def slide_12(prs):
    slide = new_slide(prs, BG_ALT)
    add_text(
        slide,
        "FECHAMENTO",
        0.72,
        0.66,
        2.2,
        0.28,
        size=11,
        color=CYAN,
        bold=True,
        tracking=1.5,
    )
    add_text(
        slide,
        "Não existe campeão universal.",
        0.72,
        1.25,
        9.5,
        0.62,
        size=34,
        bold=True,
    )
    add_text(
        slide,
        "Existe o modelo cuja hipótese combina melhor com os dados, o horizonte e a decisão.",
        0.72,
        2.00,
        9.9,
        0.50,
        size=18,
        color=MUTED,
    )

    takeaways = [
        ("01", "REPRESENTAÇÃO", "Tabular ou sequência?", ORANGE),
        ("02", "MEMÓRIA", "Curta, longa ou multiescala?", CYAN),
        ("03", "SAÍDA", "Ponto ou distribuição?", PURPLE),
    ]
    for i, (num, title, question, color) in enumerate(takeaways):
        x = 0.72 + i * 4.02
        add_card(slide, x, 3.02, 3.65, 1.58, fill=CARD)
        add_text(slide, num, x + 0.25, 3.27, 0.42, 0.24, size=11, color=color, bold=True)
        add_text(slide, title, x + 0.78, 3.27, 2.15, 0.24, size=10, color=MUTED, bold=True)
        add_text(slide, question, x + 0.25, 3.77, 3.06, 0.34, size=16, color=WHITE, bold=True)
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, 4.52, 3.65, 0.08, fill=color)

    model_colors = [ORANGE, BLUE, CYAN, LIME, PURPLE, PINK, YELLOW]
    model_names = ["XGB", "MLP", "RNN", "LSTM", "D-RNN", "DeepAR", "NPTS*"]
    for i, (name, color) in enumerate(zip(model_names, model_colors)):
        x = 0.82 + i * 1.72
        add_shape(slide, MSO_SHAPE.OVAL, x, 5.35, 0.18, 0.18, fill=color)
        add_text(slide, name, x + 0.28, 5.30, 1.15, 0.24, size=10.5, color=WHITE, bold=True)
    add_line(slide, 0.82, 5.83, 12.50, 5.83, color=GRID, width=1)
    add_text(
        slide,
        "PERGUNTAS?",
        0.72,
        6.28,
        4.0,
        0.48,
        size=24,
        color=CYAN,
        bold=True,
    )
    add_text(
        slide,
        "Comece com uma boa baseline. Evolua quando a evidência pedir.",
        5.04,
        6.32,
        7.20,
        0.34,
        size=14,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )
    return slide


def set_core_properties(prs: Presentation):
    prs.core_properties.title = "Modelos de Machine Learning"
    prs.core_properties.subject = (
        "Share Knowledge sobre XGBoost, MLP, RNN, LSTM, DilatedRNN, "
        "DeepAR e aproximação inspirada no DeepNPTS"
    )
    prs.core_properties.author = "Equipe do projeto"
    prs.core_properties.keywords = (
        "machine learning, séries temporais, XGBoost, MLP, RNN, LSTM, "
        "DilatedRNN, DeepAR, DeepNPTS"
    )


def validate(prs: Presentation):
    assert len(prs.slides) == 12
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            assert shape.left >= 0, (slide_idx, shape.name, "left")
            assert shape.top >= 0 or slide_idx == 1, (slide_idx, shape.name, "top")
            assert (
                shape.left + shape.width <= prs.slide_width + Inches(0.01)
                or (slide_idx == 1 and shape.shape_type == 1)
            ), (slide_idx, shape.name, "right")
            assert shape.top + shape.height <= prs.slide_height + Inches(0.01) or (
                slide_idx == 1 and shape.top < 0
            ), (slide_idx, shape.name, "bottom")
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                assert "\x00" not in text


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    set_core_properties(prs)

    slide_01(prs)
    slide_02(prs)
    slide_03(prs)
    slide_04(prs)
    slide_05(prs)
    slide_06(prs)
    slide_07(prs)
    slide_08(prs)
    slide_09(prs)
    slide_10(prs)
    slide_11(prs)
    slide_12(prs)

    validate(prs)
    prs.save(OUTPUT)
    print(f"Apresentação criada: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Tamanho: {OUTPUT.stat().st_size / 1024:.1f} KiB")


if __name__ == "__main__":
    main()
